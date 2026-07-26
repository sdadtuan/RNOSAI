#!/usr/bin/env python3
"""Generate PowerPoint: Bàn giao & đào tạo PTTADS Agency Operating Platform."""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
OUTPUT = ROOT / "docs" / "handover" / "PTTADS_Ban_Giao_Dao_Tao.pptx"
LOGO = ROOT / "static" / "images" / "ptt-logo.png"

NAVY = RGBColor(0x0F, 0x17, 0x2A)
GREEN = RGBColor(0x39, 0x8B, 0x43)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0x64, 0x74, 0x8B)


def _blank(prs: Presentation):
    return prs.slide_layouts[6]


def _header(slide, title: str, subtitle: str = ""):
    bar = slide.shapes.add_shape(1, 0, 0, Inches(10), Inches(0.95))
    bar.fill.solid()
    bar.fill.fore_color.rgb = NAVY
    bar.line.fill.background()
    if LOGO.is_file():
        slide.shapes.add_picture(str(LOGO), Inches(8.35), Inches(0.08), height=Inches(0.88))
    tb = slide.shapes.add_textbox(Inches(0.45), Inches(0.16), Inches(7.7), Inches(0.55))
    p = tb.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = WHITE
    if subtitle:
        sb = slide.shapes.add_textbox(Inches(0.45), Inches(0.68), Inches(7.7), Inches(0.32))
        sp = sb.text_frame.paragraphs[0]
        sp.text = subtitle
        sp.font.size = Pt(11)
        sp.font.color.rgb = GRAY


def _bullets(slide, items: list[str], top: float = 1.2):
    box = slide.shapes.add_textbox(Inches(0.55), Inches(top), Inches(8.8), Inches(5.5))
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.font.size = Pt(16)
        p.space_after = Pt(8)


def main() -> None:
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    slides: list[tuple[str, str, list[str]]] = [
        (
            "Bàn giao hệ thống PTTADS",
            "Agency Operating Platform · Nest + Next.js · 2026",
            [
                "PTT Advertising Solutions",
                "Staff: ops.pttads.vn · Client: portal.pttads.vn",
                "Phân hệ: CRM · Meta · SEO/AEO · Email Marketing",
                "Tài liệu: docs/handover/ (01–06)",
            ],
        ),
        (
            "Tổng quan kiến trúc",
            "Flask HTTP retired — canonical Nest + ops-web + portal-web",
            [
                "Nginx TLS → ops-web :3200 (staff UI)",
                "portal-web :3100 (client UI)",
                "Nest ptt-crm-api :3000 (REST + webhooks)",
                "ptt_worker + PostgreSQL + job queue",
            ],
        ),
        (
            "Staff console — CRM & Agency",
            "ops.pttads.vn/login",
            [
                "Leads · Customers · Hub HĐ · Service delivery",
                "Launch QA · Creative Hub · Campaign Write",
                "Workflow: Onboard → Deliver → Handover → Retain",
                "Agency: channel accounts, ingest, KPI",
            ],
        ),
        (
            "Meta Enterprise Ops",
            "/meta/facebook-ads · tracking · ads-ops · intelligence",
            [
                "Hub: spend, leads CRM, CPL, ROAS (T-1)",
                "Tracking: CAPI, pixel test, webhook health",
                "Ads Ops: Launch/Edit wizard + governance",
                "Portal: portal.pttads.vn/meta",
            ],
        ),
        (
            "SEO/AEO Enterprise Ops",
            "/seo/hub · clients · content · technical · reports",
            [
                "Hub executive + client health drill-down",
                "Content pipeline 13 stage + governance publish",
                "GSC/GA4 OAuth sync · AEO console",
                "Portal SEO + PDF reports",
            ],
        ),
        (
            "Email Marketing Enterprise Ops",
            "/email/hub · segments · campaigns · deliverability",
            [
                "Consent-first · suppression master",
                "Segment RFM/Lifecycle/Behavior · preflight QA",
                "Domain onboarding wizard · Grafana BI",
                "Portal approve campaign (flag pilot)",
            ],
        ),
        (
            "Client Portal",
            "portal.pttads.vn — viewer & approver",
            [
                "Dashboard KPI theo dịch vụ đang chạy",
                "Meta performance + CSV export",
                "SEO / Email approvals inbox",
                "Không truy cập cross-client",
            ],
        ),
        (
            "Phân quyền & bảo mật",
            "RBAC section keys · JWT · webhook verify",
            [
                "Staff caps theo module (crm, seo, email, meta)",
                "Portal scoped client_id",
                "Mật khẩu vault — không email plain text",
                "HTTPS only · PII redact logs",
            ],
        ),
        (
            "Nghiệm thu & hỗ trợ",
            "UAT smoke · sign-off A4 · SLA tier",
            [
                "Checklist: docs/handover/06 + form A4 in",
                "Hypercare 2–4 tuần post go-live",
                "P1: portal/webhook down — 30min ack",
                "Liên hệ AM + Tech Lead (điền tại bàn giao)",
            ],
        ),
    ]

    for title, subtitle, bullets in slides:
        slide = prs.slides.add_slide(_blank(prs))
        _header(slide, title, subtitle)
        _bullets(slide, bullets)

    OUTPUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUTPUT))
    print(f"Wrote {OUTPUT}")


if __name__ == "__main__":
    main()
