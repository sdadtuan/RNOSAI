#!/usr/bin/env python3
"""Generate PowerPoint: HR — Sơ đồ luồng & bàn giao khách hàng (có hình flowchart)."""
from __future__ import annotations

from dataclasses import dataclass
from pathlib import Path

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
OUTPUT = ROOT / "docs" / "huong-dan-su-dung" / "HR_Ban_Giao_Luu_Do.pptx"
ASSETS = ROOT / "docs" / "huong-dan-su-dung" / "assets" / "hr-handover-pptx"
LOGO = ROOT / "static" / "images" / "ptt-logo.png"
FONT_PATH = "/System/Library/Fonts/Supplemental/Arial Unicode.ttf"
FALLBACK_FONT = "/Library/Fonts/Arial.ttf"

NAVY = (15, 23, 42)
GREEN = (74, 124, 92)
GREEN_LIGHT = (232, 245, 236)
BLUE = (59, 130, 246)
BLUE_LIGHT = (219, 234, 254)
PURPLE = (99, 102, 241)
PURPLE_LIGHT = (224, 231, 255)
ORANGE = (234, 88, 12)
ORANGE_LIGHT = (255, 237, 213)
GRAY = (100, 116, 139)
GRAY_LIGHT = (241, 245, 249)
WHITE = (255, 255, 255)
RED_LIGHT = (254, 226, 226)


@dataclass
class Box:
    x: int
    y: int
    w: int
    h: int
    text: str
    fill: tuple[int, int, int] = GREEN_LIGHT
    border: tuple[int, int, int] = GREEN
    text_color: tuple[int, int, int] = NAVY
    radius: int = 12
    font_size: int = 15


@dataclass
class Arrow:
    x1: int
    y1: int
    x2: int
    y2: int
    color: tuple[int, int, int] = GRAY


class DiagramCanvas:
    def __init__(self, width: int = 1600, height: int = 900, title: str = ""):
        self.img = Image.new("RGB", (width, height), WHITE)
        self.draw = ImageDraw.Draw(self.img)
        self.w, self.h = width, height
        self.title = title
        self._title_font = self._font(26, bold=True)
        self._label_font = self._font(18, bold=True)
        self._small_font = self._font(13)

    def _font(self, size: int, bold: bool = False) -> ImageFont.FreeTypeFont | ImageFont.ImageFont:
        for path in (FONT_PATH, FALLBACK_FONT):
            try:
                return ImageFont.truetype(path, size)
            except OSError:
                continue
        return ImageFont.load_default()

    def header_band(self):
        self.draw.rectangle([0, 0, self.w, 72], fill=NAVY)
        if self.title:
            self.draw.text((36, 20), self.title, fill=WHITE, font=self._title_font)

    def section_label(self, text: str, x: int, y: int, color: tuple[int, int, int] = GREEN):
        tw = self.draw.textlength(text, font=self._label_font)
        pad_x, pad_y = 14, 6
        self.draw.rounded_rectangle([x, y, x + tw + pad_x * 2, y + 34], radius=8, fill=color)
        self.draw.text((x + pad_x, y + pad_y), text, fill=WHITE, font=self._label_font)

    def _wrap(self, text: str, font: ImageFont.FreeTypeFont, max_w: int) -> list[str]:
        words = text.replace("\n", " ").split()
        lines: list[str] = []
        cur = ""
        for word in words:
            test = f"{cur} {word}".strip()
            if self.draw.textlength(test, font=font) <= max_w:
                cur = test
            else:
                if cur:
                    lines.append(cur)
                cur = word
        if cur:
            lines.append(cur)
        return lines or [text]

    def box(self, b: Box):
        self.draw.rounded_rectangle(
            [b.x, b.y, b.x + b.w, b.y + b.h], radius=b.radius, fill=b.fill, outline=b.border, width=2
        )
        font = self._font(b.font_size)
        lines: list[str] = []
        for part in b.text.split("\n"):
            lines.extend(self._wrap(part, font, b.w - 20))
        line_h = b.font_size + 6
        total_h = len(lines) * line_h
        start_y = b.y + (b.h - total_h) // 2
        for i, line in enumerate(lines):
            tw = self.draw.textlength(line, font=font)
            self.draw.text((b.x + (b.w - tw) // 2, start_y + i * line_h), line, fill=b.text_color, font=font)

    def diamond(self, cx: int, cy: int, size: int, text: str, fill: tuple[int, int, int] = ORANGE_LIGHT, border: tuple[int, int, int] = ORANGE):
        pts = [(cx, cy - size), (cx + size, cy), (cx, cy + size), (cx - size, cy)]
        self.draw.polygon(pts, fill=fill, outline=border)
        font = self._font(14, bold=True)
        lines = self._wrap(text, font, size * 2 - 10)
        start_y = cy - (len(lines) * 18) // 2
        for i, line in enumerate(lines):
            tw = self.draw.textlength(line, font=font)
            self.draw.text((cx - tw // 2, start_y + i * 18), line, fill=NAVY, font=font)

    def arrow(self, a: Arrow, width: int = 3):
        import math

        self.draw.line([a.x1, a.y1, a.x2, a.y2], fill=a.color, width=width)
        angle = math.atan2(a.y2 - a.y1, a.x2 - a.x1)
        head = 12
        left = (a.x2 - head * math.cos(angle - math.pi / 6), a.y2 - head * math.sin(angle - math.pi / 6))
        right = (a.x2 - head * math.cos(angle + math.pi / 6), a.y2 - head * math.sin(angle + math.pi / 6))
        self.draw.polygon([(a.x2, a.y2), left, right], fill=a.color)

    def save(self, path: Path):
        path.parent.mkdir(parents=True, exist_ok=True)
        self.img.save(path, "PNG")


def diagram_module_map(path: Path):
    c = DiagramCanvas(title="§1 — Bản đồ module HR trên RNOSAI")
    c.header_band()
    c.section_label("Nền tảng (trước HR)", 40, 95, NAVY)
    for i, t in enumerate(["Tạo user\n/admin/org/users/new", "RBAC\n/permissions", "Org chart\n/admin/crm/org"]):
        c.box(Box(40 + i * 390, 140, 350, 72, t, fill=GRAY_LIGHT, border=GRAY))
    c.box(Box(640, 240, 320, 62, "HR Hub\n/crm/hr", fill=GREEN, border=GREEN, text_color=WHITE, font_size=17))
    for i in range(3):
        c.arrow(Arrow(215 + i * 390, 212, 800, 240))
    c.section_label("Employee File P1–P5", 40, 330, GREEN)
    tabs = ["P1 Hồ sơ\nĐịnh danh", "P4 Ví giấy tờ", "P2 Hợp đồng", "P3 Bảo hiểm", "P5 Gia đình\nLifecycle"]
    for i, t in enumerate(tabs):
        c.box(Box(40 + i * 310, 380, 280, 72, t))
        if i < 4:
            c.arrow(Arrow(320 + i * 310, 416, 350 + i * 310, 416))
    c.arrow(Arrow(800, 302, 800, 380))
    c.section_label("Self-service NV", 40, 480, BLUE)
    for i, t in enumerate(["Ví tự nộp P6\n/my-wallet", "GPS P8\n/payroll/me", "Nghỉ phép\n/hr/leave", "Phiếu lương\n/payroll/me"]):
        c.box(Box(40 + i * 390, 530, 350, 72, t, fill=BLUE_LIGHT, border=BLUE))
    c.section_label("Chấm công P7–P8", 40, 630, PURPLE)
    att = ["Trung tâm\n/hr/attendance", "Máy ZK + PIN", "Site GPS", "Rollup ngày"]
    for i, t in enumerate(att):
        c.box(Box(40 + i * 390, 680, 350, 62, t, fill=PURPLE_LIGHT, border=PURPLE))
        if i < 3:
            c.arrow(Arrow(390 + i * 390, 711, 430 + i * 390, 711, color=PURPLE))
    c.box(Box(1200, 680, 360, 62, "Payroll lite\n/crm/payroll", fill=GREEN, border=GREEN, text_color=WHITE))
    c.arrow(Arrow(1600, 711, 1200, 711, color=GREEN))
    c.save(path)


def diagram_e2e_phases(path: Path):
    c = DiagramCanvas(title="§2 — Luồng chính: NV mới → Vận hành → Offboard")
    c.header_band()
    phases = [
        ("0\nSetup", "Flag HR=1\nDDL P1–P8\nRBAC + Org", GRAY_LIGHT, GRAY),
        ("1\nNV mới", "Tạo user CRM\nRoster /crm/staff\nOffer → Onboard", GREEN_LIGHT, GREEN),
        ("2\nHồ sơ 360", "P1 Định danh\nP4 Ví · P2 HĐ\nP3 BH · P5 NPT", GREEN_LIGHT, GREEN),
        ("3\nChính thức", "Gate BR-HR-130\nLifecycle official\nHĐ chính thức", BLUE_LIGHT, BLUE),
        ("4\nVận hành", "Chấm công P7/P8\nVí P6 · Nghỉ phép\nPayroll tháng", BLUE_LIGHT, BLUE),
        ("5\nOffboard", "Thông báo nghỉ\nOffboard\nLưu trữ", ORANGE_LIGHT, ORANGE),
    ]
    y = 180
    bw, bh, gap = 240, 210, 16
    for i, (title, body, fill, border) in enumerate(phases):
        x = 30 + i * (bw + gap)
        c.box(Box(x, y, bw, 55, title, fill=border, border=border, text_color=WHITE, font_size=18))
        c.box(Box(x, y + 65, bw, bh, body, fill=fill, border=border))
        if i < len(phases) - 1:
            c.arrow(Arrow(x + bw, y + 155, x + bw + gap, y + 155))
    c.diamond(800, 520, 75, "Gate BR-HR-130\nHĐ+CCCD+Địa chỉ?")
    c.box(Box(480, 620, 280, 58, "Thiếu → quay P1\nbổ sung hồ sơ", fill=RED_LIGHT, border=ORANGE))
    c.box(Box(920, 620, 280, 58, "Đủ → Chính thức\n→ Vận hành", fill=GREEN_LIGHT, border=GREEN))
    c.arrow(Arrow(725, 595, 620, 620, color=ORANGE))
    c.arrow(Arrow(875, 595, 1060, 620, color=GREEN))
    c.save(path)


def diagram_roles(path: Path):
    c = DiagramCanvas(title="§3 — Luồng theo vai trò (ai làm gì hàng ngày)")
    c.header_band()
    lanes = [
        ("HR / HCNS", GREEN, ["HR Hub cảnh báo", "Duyệt ví P6", "Duyệt GPS P8", "Cập nhật hồ sơ 360"]),
        ("Quản lý trực tiếp", BLUE, ["Duyệt đơn nghỉ phép", "Xem roster team\n(mask PII)"]),
        ("Nhân viên", PURPLE, ["Chấm GPS vào/ra", "Nộp my-wallet", "Xin nghỉ phép", "Xem phiếu lương"]),
        ("IT / Vận hành", GRAY, ["Thiết bị ADMS", "Import CSV chấm công", "Site geofence GPS"]),
        ("Kế toán", ORANGE, ["Export Excel ví+NPT", "Review payroll tháng"]),
    ]
    y = 110
    for name, color, steps in lanes:
        c.section_label(name, 30, y, color)
        for i, step in enumerate(steps):
            fill = GRAY_LIGHT if color == GRAY else GREEN_LIGHT
            if color == BLUE:
                fill = BLUE_LIGHT
            elif color == PURPLE:
                fill = PURPLE_LIGHT
            elif color == ORANGE:
                fill = ORANGE_LIGHT
            c.box(Box(30 + i * 390, y + 42, 360, 62, step, fill=fill, border=color))
            if i < len(steps) - 1:
                c.arrow(Arrow(390 + i * 390, y + 73, 420 + i * 390, y + 73))
        y += 130
    # cross-role links
    c.draw.text((30, 780), "Liên kết: NV nộp ví → HR duyệt · NV GPS → HR duyệt ngoài vùng · NV xin nghỉ → Manager duyệt · Rollup → Kế toán payroll",
                fill=GRAY, font=c._small_font)
    c.save(path)


def diagram_attendance_payroll(path: Path):
    c = DiagramCanvas(title="§4 — Luồng dữ liệu: Chấm công → Lương")
    c.header_band()
    c.box(Box(40, 150, 260, 72, "Máy ZK/ADMS\nPIN", fill=PURPLE_LIGHT, border=PURPLE))
    c.box(Box(40, 280, 260, 72, "NV PWA GPS\n/crm/payroll/me", fill=BLUE_LIGHT, border=BLUE))
    c.box(Box(380, 200, 300, 80, "hr_attendance_punches\nsource=device / gps", fill=GRAY_LIGHT, border=GRAY))
    c.arrow(Arrow(300, 186, 380, 220))
    c.arrow(Arrow(300, 316, 380, 260))
    c.diamond(780, 240, 70, "pending\nreview?")
    c.arrow(Arrow(680, 240, 710, 240))
    c.box(Box(920, 150, 240, 62, "accepted", fill=GREEN_LIGHT, border=GREEN))
    c.box(Box(920, 320, 240, 62, "rejected", fill=RED_LIGHT, border=ORANGE))
    c.arrow(Arrow(850, 210, 920, 180, color=GREEN))
    c.arrow(Arrow(850, 270, 920, 350, color=ORANGE))
    c.draw.text((860, 120), "HR duyệt", fill=GREEN, font=c._small_font)
    c.draw.text((860, 300), "Từ chối", fill=ORANGE, font=c._small_font)
    c.box(Box(1220, 150, 340, 72, "Rollup TZ VN\nBR-HR-154: máy thắng", fill=GREEN_LIGHT, border=GREEN))
    c.arrow(Arrow(1160, 181, 1220, 186))
    c.box(Box(1220, 280, 340, 62, "crm_attendance\n1 dòng/ngày/NV", fill=BLUE_LIGHT, border=BLUE))
    c.arrow(Arrow(1390, 222, 1390, 280))
    c.box(Box(880, 430, 280, 62, "Nghỉ phép\n/crm/hr/leave", fill=PURPLE_LIGHT, border=PURPLE))
    c.box(Box(1220, 420, 340, 72, "Payroll lite\n/crm/payroll", fill=GREEN, border=GREEN, text_color=WHITE))
    c.arrow(Arrow(1390, 342, 1390, 420))
    c.arrow(Arrow(1160, 461, 1220, 456, color=PURPLE))
    c.box(Box(1220, 540, 340, 62, "Phiếu lương NV\n/crm/payroll/me", fill=BLUE_LIGHT, border=BLUE))
    c.arrow(Arrow(1390, 492, 1390, 540))
    c.save(path)


def diagram_lifecycle(path: Path):
    c = DiagramCanvas(title="§5 — Lifecycle 8 stage (P5) — mốc bàn giao")
    c.header_band()
    stages = [
        ("Offer", "Đồng ý tuyển\nTạo user + org", GRAY_LIGHT, GRAY),
        ("Onboard\ngiấy tờ", "Thu hồ sơ\nP1 + P4", GREEN_LIGHT, GREEN),
        ("Thử việc", "HĐ thử việc P2\nChấm công", GREEN_LIGHT, GREEN),
        ("Chính thức", "Gate: HĐ+CCCD\n+ địa chỉ", BLUE_LIGHT, BLUE),
        ("Chuyển BP", "Đổi phòng/team\nPhụ lục HĐ", BLUE_LIGHT, BLUE),
        ("TB nghỉ", "Báo nghỉ\nHandover", ORANGE_LIGHT, ORANGE),
        ("Offboard", "Thủ tục nghỉ\nQuyết toán", ORANGE_LIGHT, ORANGE),
        ("Lưu trữ", "NV inactive\nChỉ xem", GRAY_LIGHT, GRAY),
    ]
    for i, (title, body, fill, border) in enumerate(stages):
        row, col = i // 4, i % 4
        x, y = 40 + col * 390, 140 + row * 200
        c.box(Box(x, y, 80, 200 if row == 0 else 180, title, fill=border, border=border, text_color=WHITE, font_size=16))
        c.box(Box(x + 90, y, 280, 200 if row == 0 else 180, body, fill=fill, border=border))
        if col < 3 and row == 0:
            c.arrow(Arrow(x + 370, y + 90, x + 390, y + 90))
        if i == 3:
            c.arrow(Arrow(x + 370, y + 90, 40, y + 290))
        if col < 3 and row == 1:
            c.arrow(Arrow(x + 370, y + 90, x + 390, y + 90))
    c.box(Box(40, 560, 1520, 58, "Gate BR-HR-130: Chuyển Chính thức chỉ khi có HĐ active + CCCD + địa chỉ thường trú", fill=ORANGE_LIGHT, border=ORANGE, font_size=16))
    c.save(path)


def diagram_one_page(path: Path):
    c = DiagramCanvas(title="Sơ đồ tổng hợp — NV mới → Offboard (training)")
    c.header_band()
    c.box(Box(680, 100, 240, 58, "NV MỚI", fill=GREEN, border=GREEN, text_color=WHITE, font_size=17))
    c.arrow(Arrow(800, 158, 800, 190))
    c.box(Box(560, 200, 480, 58, "Setup: Flag + DDL + Org + RBAC", fill=GRAY_LIGHT, border=GRAY))
    c.arrow(Arrow(800, 258, 800, 290))
    c.box(Box(560, 300, 480, 58, "Hồ sơ 360: P1 → P4 → P2 → P3 → P5", fill=GREEN_LIGHT, border=GREEN))
    c.diamond(800, 420, 65, "Gate\nChính thức?")
    c.box(Box(320, 500, 300, 58, "Bổ sung\nhồ sơ P1/P4", fill=RED_LIGHT, border=ORANGE))
    c.box(Box(980, 500, 300, 58, "Vận hành:\nChấm công · Leave · Payroll", fill=BLUE_LIGHT, border=BLUE))
    c.arrow(Arrow(735, 485, 470, 500, color=ORANGE))
    c.arrow(Arrow(865, 485, 1130, 500, color=BLUE))
    c.box(Box(560, 620, 480, 58, "Self-service: my-wallet · GPS · payslip", fill=PURPLE_LIGHT, border=PURPLE))
    c.arrow(Arrow(1130, 558, 900, 620, color=PURPLE))
    c.box(Box(560, 720, 480, 58, "Offboard → Lưu trữ", fill=ORANGE_LIGHT, border=ORANGE))
    c.arrow(Arrow(800, 678, 800, 720))
    c.save(path)


def generate_diagrams() -> dict[str, Path]:
    ASSETS.mkdir(parents=True, exist_ok=True)
    specs = {
        "01_module_map": diagram_module_map,
        "02_e2e_phases": diagram_e2e_phases,
        "03_roles": diagram_roles,
        "04_attendance_payroll": diagram_attendance_payroll,
        "05_lifecycle": diagram_lifecycle,
        "06_one_page": diagram_one_page,
    }
    paths: dict[str, Path] = {}
    for key, fn in specs.items():
        p = ASSETS / f"diagram_{key}.png"
        fn(p)
        paths[key] = p
    return paths


PPTX_NAVY = RGBColor(0x0F, 0x17, 0x2A)
PPTX_GREEN = RGBColor(0x4A, 0x7C, 0x5C)
PPTX_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
PPTX_GRAY = RGBColor(0x64, 0x74, 0x8B)


def _blank(prs: Presentation):
    return prs.slide_layouts[6]


def _header(slide, title: str, subtitle: str = ""):
    bar = slide.shapes.add_shape(1, 0, 0, Inches(10), Inches(0.95))
    bar.fill.solid()
    bar.fill.fore_color.rgb = PPTX_NAVY
    bar.line.fill.background()
    if LOGO.is_file():
        slide.shapes.add_picture(str(LOGO), Inches(8.35), Inches(0.08), height=Inches(0.88))
    tb = slide.shapes.add_textbox(Inches(0.45), Inches(0.16), Inches(7.7), Inches(0.55))
    p = tb.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = PPTX_WHITE
    if subtitle:
        sb = slide.shapes.add_textbox(Inches(0.45), Inches(0.68), Inches(7.7), Inches(0.32))
        sp = sb.text_frame.paragraphs[0]
        sp.text = subtitle
        sp.font.size = Pt(11)
        sp.font.color.rgb = PPTX_GRAY


def _bullets(slide, items: list[str], top: float = 1.2, size: int = 16):
    box = slide.shapes.add_textbox(Inches(0.55), Inches(top), Inches(8.8), Inches(5.8))
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.font.size = Pt(size)
        p.space_after = Pt(8)


def _diagram_slide(prs: Presentation, title: str, subtitle: str, img: Path):
    slide = prs.slides.add_slide(_blank(prs))
    _header(slide, title, subtitle)
    slide.shapes.add_picture(str(img), Inches(0.35), Inches(1.05), width=Inches(9.3))


def main() -> None:
    diagrams = generate_diagrams()

    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Title slide
    slide = prs.slides.add_slide(_blank(prs))
    bar = slide.shapes.add_shape(1, 0, 0, Inches(10), Inches(7.5))
    bar.fill.solid()
    bar.fill.fore_color.rgb = PPTX_NAVY
    bar.line.fill.background()
    if LOGO.is_file():
        slide.shapes.add_picture(str(LOGO), Inches(3.8), Inches(0.8), height=Inches(1.2))
    tb = slide.shapes.add_textbox(Inches(0.6), Inches(2.3), Inches(8.8), Inches(2.5))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = "HR — Sơ đồ luồng quản lý"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = PPTX_WHITE
    p2 = tf.add_paragraph()
    p2.text = "Khung bàn giao & đào tạo khách hàng"
    p2.font.size = Pt(22)
    p2.font.color.rgb = PPTX_GREEN
    p3 = tf.add_paragraph()
    p3.text = "RNOSAI · rs.pttads.vn · Phiên bản 1.0 · 2026-08-25"
    p3.font.size = Pt(14)
    p3.font.color.rgb = PPTX_GRAY

    slide = prs.slides.add_slide(_blank(prs))
    _header(slide, "Mục tiêu bàn giao", "Tài liệu: docs/huong-dan-su-dung/22-hr-handover-flow-and-guides.md")
    _bullets(
        slide,
        [
            "Sơ đồ end-to-end: NV mới → Hồ sơ 360 → Chính thức → Vận hành → Offboard",
            "Employee File OS P1–P8: hồ sơ, ví, HĐ, BH, chấm công, payroll",
            "Vai trò: HR · Manager · NV · IT · Kế toán",
            "18 file hướng dẫn chi tiết (22a–22r) — soạn theo template bàn giao",
            "Kịch bản training tuần đầu cho NV mới",
        ],
    )

    for title, subtitle, key in [
        ("Bản đồ module HR", "Platform → HR Hub → Employee File → Chấm công → Payroll", "01_module_map"),
        ("Luồng chính end-to-end", "Giai đoạn 0 (Setup) → 5 (Offboard) + Gate BR-HR-130", "02_e2e_phases"),
        ("Luồng theo vai trò", "HR · Manager · NV · IT · Kế toán", "03_roles"),
        ("Luồng chấm công → lương", "Máy/GPS → punch → rollup → payroll → payslip", "04_attendance_payroll"),
        ("Lifecycle 8 stage", "Offer → Onboard → Thử việc → Chính thức → … → Lưu trữ", "05_lifecycle"),
        ("Sơ đồ tổng hợp một trang", "Poster training / wall chart", "06_one_page"),
    ]:
        _diagram_slide(prs, title, subtitle, diagrams[key])

    slide = prs.slides.add_slide(_blank(prs))
    _header(slide, "Kịch bản tuần đầu — NV mới", "§7 — Training HR + NV mới")
    _bullets(
        slide,
        [
            "D0: HR tạo user + gán org · NV đăng nhập lần đầu → thấy menu HR",
            "D1: HR nhập định danh + địa chỉ (P1) · Tạo thẻ CCCD + upload ví (P4)",
            "D2: HR tạo HĐ thử việc (P2) · Gán PIN máy chấm công (nếu dùng máy)",
            "D3: HR tạo site GPS + gán NV · NV chấm vào ca trên /crm/payroll/me",
            "D3–D5: NV nộp bằng/chứng chỉ tại /crm/hr/my-wallet",
            "D5: HR Hub duyệt ví + GPS pending → ví % tăng, rollup OK",
            "Cuối tháng: HR/Kế toán review /crm/payroll · NV xem payslip",
        ],
        size=15,
    )

    slide = prs.slides.add_slide(_blank(prs))
    _header(slide, "Ma trận tính năng (tóm tắt)", "§6 — 18 tính năng · file chi tiết 22a–22r")
    _bullets(
        slide,
        [
            "P0 go-live: 22a Setup · 22b NV mới · 22e Định danh · 22f Ví HR · 22c HR Hub",
            "P1 gate chính thức: 22g HĐ · 22h BH · 22i Lifecycle · 22j Self-wallet",
            "P1 chấm công: 22l Máy chấm · 22m GPS · 22n Tab chấm công",
            "P2 cuối tháng: 22o Leave · 22p Payroll admin · 22q Payslip NV",
            "P2 admin: 22k Export kế toán · 22r RBAC matrix",
            "Tham chiếu đầy đủ: 17-hr-employee-file-os.md",
        ],
        size=15,
    )

    slide = prs.slides.add_slide(_blank(prs))
    _header(slide, "Điều kiện vận hành", "§10 — Checklist IT trước bàn giao")
    _bullets(
        slide,
        [
            "PTT_HR_EMPLOYEE_FILE=1 trong .env ptt-crm-api",
            "DDL P1–P8: scripts/apply_pg_ddl_hr_employee_file_p*.sh",
            "Ops-web build: NEXT_PUBLIC_WIN_ORG_UI=1, WIN_LEAVE_LITE=1, WIN_PAYSLIP_PORTAL=1",
            "Smoke: bash scripts/smoke_hr_employee_file_p8.sh pass",
            "RBAC: crm_staff_roster.*, crm_hr_pii.*, crm_hr_docs.*, crm_payroll_*",
        ],
        size=15,
    )

    slide = prs.slides.add_slide(_blank(prs))
    _header(slide, "Bước tiếp theo", "Soạn file bàn giao P0 theo template §8")
    _bullets(
        slide,
        [
            "1. IT hoàn tất flag + DDL + smoke test (22a)",
            "2. HR training: tạo NV + hồ sơ P1/P4 (22b, 22e, 22f)",
            "3. Pilot 1 NV qua gate Chính thức (22i) trước go-live rộng",
            "4. Bật chấm công máy/GPS theo nhu cầu (22l, 22m)",
            "5. Cuối tháng đầu: chạy payroll lite + payslip (22p, 22q)",
            "",
            "Liên hệ: HR dự án + Tech Lead PTT (điền tại buổi bàn giao)",
        ],
        size=16,
    )

    OUTPUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUTPUT))
    print(f"Wrote {OUTPUT}")
    print(f"Diagram PNGs: {ASSETS}/")


if __name__ == "__main__":
    main()
