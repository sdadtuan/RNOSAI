#!/usr/bin/env python3
"""Generate PowerPoint: Leads — Sơ đồ luồng & bàn giao khách hàng (có hình flowchart)."""
from __future__ import annotations

import textwrap
from dataclasses import dataclass
from pathlib import Path

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
OUTPUT = ROOT / "docs" / "huong-dan-su-dung" / "Leads_Ban_Giao_Luu_Do.pptx"
ASSETS = ROOT / "docs" / "huong-dan-su-dung" / "assets" / "leads-handover-pptx"
LOGO = ROOT / "static" / "images" / "ptt-logo.png"
FONT_PATH = "/System/Library/Fonts/Supplemental/Arial Unicode.ttf"
FALLBACK_FONT = "/Library/Fonts/Arial.ttf"

# Brand palette (aligned with ops-web sage theme)
NAVY = (15, 23, 42)
GREEN = (74, 124, 92)
GREEN_LIGHT = (232, 245, 236)
BLUE = (59, 130, 246)
BLUE_LIGHT = (219, 234, 254)
ORANGE = (234, 88, 12)
ORANGE_LIGHT = (255, 237, 213)
GRAY = (100, 116, 139)
GRAY_LIGHT = (241, 245, 249)
WHITE = (255, 255, 255)
RED_LIGHT = (254, 226, 226)


@dataclass
class Box:
    x: int
    y: int
    w: int
    h: int
    text: str
    fill: tuple[int, int, int] = GREEN_LIGHT
    border: tuple[int, int, int] = GREEN
    text_color: tuple[int, int, int] = NAVY
    radius: int = 12
    font_size: int = 15


@dataclass
class Arrow:
    x1: int
    y1: int
    x2: int
    y2: int
    color: tuple[int, int, int] = GRAY


class DiagramCanvas:
    def __init__(self, width: int = 1600, height: int = 900, title: str = ""):
        self.img = Image.new("RGB", (width, height), WHITE)
        self.draw = ImageDraw.Draw(self.img)
        self.w, self.h = width, height
        self.title = title
        self._title_font = self._font(26, bold=True)
        self._box_font = self._font(15)
        self._small_font = self._font(13)
        self._label_font = self._font(18, bold=True)

    def _font(self, size: int, bold: bool = False) -> ImageFont.FreeTypeFont | ImageFont.ImageFont:
        for path in (FONT_PATH, FALLBACK_FONT):
            try:
                return ImageFont.truetype(path, size)
            except OSError:
                continue
        return ImageFont.load_default()

    def header_band(self):
        self.draw.rectangle([0, 0, self.w, 72], fill=NAVY)
        if self.title:
            self.draw.text((36, 20), self.title, fill=WHITE, font=self._title_font)

    def section_label(self, text: str, x: int, y: int, color: tuple[int, int, int] = GREEN):
        tw = self.draw.textlength(text, font=self._label_font)
        pad_x, pad_y = 14, 6
        self.draw.rounded_rectangle(
            [x, y, x + tw + pad_x * 2, y + 34],
            radius=8,
            fill=color,
        )
        self.draw.text((x + pad_x, y + pad_y), text, fill=WHITE, font=self._label_font)

    def _wrap(self, text: str, font: ImageFont.FreeTypeFont, max_w: int) -> list[str]:
        words = text.replace("\n", " ").split()
        lines: list[str] = []
        cur = ""
        for word in words:
            test = f"{cur} {word}".strip()
            if self.draw.textlength(test, font=font) <= max_w:
                cur = test
            else:
                if cur:
                    lines.append(cur)
                cur = word
        if cur:
            lines.append(cur)
        return lines or [text]

    def box(self, b: Box):
        self.draw.rounded_rectangle([b.x, b.y, b.x + b.w, b.y + b.h], radius=b.radius, fill=b.fill, outline=b.border, width=2)
        font = self._font(b.font_size)
        lines = []
        for part in b.text.split("\n"):
            lines.extend(self._wrap(part, font, b.w - 20))
        line_h = b.font_size + 6
        total_h = len(lines) * line_h
        start_y = b.y + (b.h - total_h) // 2
        for i, line in enumerate(lines):
            tw = self.draw.textlength(line, font=font)
            tx = b.x + (b.w - tw) // 2
            ty = start_y + i * line_h
            self.draw.text((tx, ty), line, fill=b.text_color, font=font)

    def diamond(self, cx: int, cy: int, size: int, text: str, fill: tuple[int, int, int] = ORANGE_LIGHT, border: tuple[int, int, int] = ORANGE):
        pts = [(cx, cy - size), (cx + size, cy), (cx, cy + size), (cx - size, cy)]
        self.draw.polygon(pts, fill=fill, outline=border)
        font = self._font(14, bold=True)
        lines = self._wrap(text, font, size * 2 - 10)
        line_h = 18
        start_y = cy - (len(lines) * line_h) // 2
        for i, line in enumerate(lines):
            tw = self.draw.textlength(line, font=font)
            self.draw.text((cx - tw // 2, start_y + i * line_h), line, fill=NAVY, font=font)

    def arrow(self, a: Arrow, width: int = 3):
        self.draw.line([a.x1, a.y1, a.x2, a.y2], fill=a.color, width=width)
        # arrowhead
        import math

        angle = math.atan2(a.y2 - a.y1, a.x2 - a.x1)
        head = 12
        left = (
            a.x2 - head * math.cos(angle - math.pi / 6),
            a.y2 - head * math.sin(angle - math.pi / 6),
        )
        right = (
            a.x2 - head * math.cos(angle + math.pi / 6),
            a.y2 - head * math.sin(angle + math.pi / 6),
        )
        self.draw.polygon([(a.x2, a.y2), left, right], fill=a.color)

    def save(self, path: Path):
        path.parent.mkdir(parents=True, exist_ok=True)
        self.img.save(path, "PNG")


def diagram_two_flows(path: Path):
    c = DiagramCanvas(title="§0 — Hai luồng Lead (đừng nhầm khi bàn giao)")
    c.header_band()
    # B2B column
    c.section_label("B2B Sales — Prospect mới", 60, 100, GREEN)
    for i, txt in enumerate(
        [
            "Mục đích: Bán HĐ agency MỚI",
            "Danh sách: /crm/b2b/leads",
            "Bắt buộc: Dự án PTT",
            "Thắng: status WON",
            "SLA: First touch 5–15 phút",
        ]
    ):
        c.box(Box(60, 150 + i * 95, 680, 78, txt, fill=GREEN_LIGHT))
    # CSKH column
    c.section_label("CSKH vận hành — Client spa", 860, 100, BLUE)
    for i, txt in enumerate(
        [
            "Mục đích: Lead ads KH đã ký HĐ",
            "Danh sách: /crm/operational/leads",
            "Bắt buộc: Khách hàng agency",
            "Thắng: status CHỐT",
            "SLA: Meta 24h (SOP spa)",
        ]
    ):
        c.box(Box(860, 150 + i * 95, 680, 78, txt, fill=BLUE_LIGHT))
    c.box(
        Box(
            420,
            720,
            760,
            70,
            "⚠ Quy tắc vàng: Prospect B2B ≠ lead spa của khách đang chạy ads",
            fill=ORANGE_LIGHT,
            border=ORANGE,
            font_size=17,
        )
    )
    c.save(path)


def diagram_module_map(path: Path):
    c = DiagramCanvas(title="§1 — Bản đồ module Leads trên RNOSAI")
    c.header_band()
    c.section_label("Setup IT + GDKD", 40, 95, NAVY)
    setup = ["DDL + Flag B2B OS", "Dự án PTT\n+ map kênh", "Pool NV\n+ SLA", "Nguồn & Kênh"]
    for i, t in enumerate(setup):
        c.box(Box(40 + i * 390, 140, 350, 70, t, fill=GRAY_LIGHT, border=GRAY))
    c.section_label("Nguồn lead", 40, 240, GREEN)
    sources = ["Facebook", "Zalo OA", "Google", "Webform", "API", "Nhập tay"]
    for i, t in enumerate(sources):
        c.box(Box(40 + (i % 3) * 390, 290 + (i // 3) * 80, 350, 62, t))
    c.box(Box(1220, 300, 340, 62, "Ingress\nchưa map", fill=RED_LIGHT, border=ORANGE))
    c.arrow(Arrow(780, 360, 1220, 330))
    c.section_label("Pipeline CRM", 40, 430, GREEN)
    pipe = ["Tạo lead\n+ AI score", "Auto-assign\npool dự án", "Inbox B2B\nalert Hot", "Chi tiết lead\n+ Gọi"]
    for i, t in enumerate(pipe):
        c.box(Box(40 + i * 390, 480, 350, 78, t))
        if i < 3:
            c.arrow(Arrow(390 + i * 390, 519, 430 + i * 390, 519))
    c.section_label("Pre-sales → Won", 40, 590, BLUE)
    ps = ["Intake BANT", "Solution\nqueue", "Deal Room\n+ R5", "Proposal\n+ Hub HĐ", "Delivery"]
    for i, t in enumerate(ps):
        c.box(Box(40 + i * 310, 640, 280, 78, t, fill=BLUE_LIGHT, border=BLUE))
        if i < 4:
            c.arrow(Arrow(320 + i * 310, 679, 350 + i * 310, 679, color=BLUE))
    c.section_label("GDKD giám sát", 40, 760, ORANGE)
    c.box(Box(40, 810, 480, 58, "Command center · Speed-to-lead · Review queue", fill=ORANGE_LIGHT, border=ORANGE))
    c.box(Box(560, 810, 480, 58, "KPI GDKD Enterprise · Bảng CSKH SLA", fill=ORANGE_LIGHT, border=ORANGE))
    c.save(path)


def diagram_e2e_phases(path: Path):
    c = DiagramCanvas(title="§2 — Luồng chính: Lead mới → Won → Delivery")
    c.header_band()
    phases = [
        ("0\nSetup", "IT + GDKD\nDDL · Dự án PTT\nPool NV", GRAY_LIGHT, GRAY),
        ("1\nLead mới", "Webhook/API\nAI score\nAuto-assign", GREEN_LIGHT, GREEN),
        ("2\nFirst touch", "Inbox alert\nGọi ≤15p\nLog Activity", GREEN_LIGHT, GREEN),
        ("3\nQualify", "Intake BANT\nReview queue\nGate B2 ✓", BLUE_LIGHT, BLUE),
        ("4\nPre-sales", "Consult\nSolution R5\nDeal Room", BLUE_LIGHT, BLUE),
        ("5\nWon", "Proposal\nHub Active\nOnboard DV", GREEN_LIGHT, GREEN),
    ]
    x0, y = 30, 180
    bw, bh, gap = 240, 200, 18
    for i, (title, body, fill, border) in enumerate(phases):
        x = x0 + i * (bw + gap)
        c.box(Box(x, y, bw, 55, title, fill=border, border=border, text_color=WHITE, font_size=18))
        c.box(Box(x, y + 65, bw, bh, body, fill=fill, border=border))
        if i < len(phases) - 1:
            c.arrow(Arrow(x + bw, y + 145, x + bw + gap, y + 145))
    # CSKH branch
    c.section_label("Nhánh CSKH spa (song song)", 30, 430, BLUE)
    spa = ["Lead spa Meta/Zalo", "Bảng CSKH SLA", "Gọi ≤24h", "Chốt / Lost"]
    for i, t in enumerate(spa):
        c.box(Box(30 + i * 390, 480, 350, 72, t, fill=BLUE_LIGHT, border=BLUE))
        if i < 3:
            c.arrow(Arrow(380 + i * 390, 516, 420 + i * 390, 516, color=BLUE))
    c.save(path)


def diagram_roles(path: Path):
    c = DiagramCanvas(title="§3 — Luồng theo vai trò (ai làm gì hàng ngày)")
    c.header_band()
    lanes = [
        ("AM / Sales B2B", GREEN, ["Inbox Hot → Gọi 15p", "Intake BANT", "Handoff Solution", "Deal Room + Won"]),
        ("CSKH vận hành", BLUE, ["Bảng SLA lead mới", "Gọi spa ≤24h", "Cập nhật chốt/lost"]),
        ("Solution", (99, 102, 241), ["Claim queue", "KH MKT sơ bộ R5", "Tham gia Deal Room 45p"]),
        ("GDKD", ORANGE, ["Command center SLA", "Review queue release", "Approve deal lớn"]),
        ("Marketing / IT", GRAY, ["Campaign Meta/Zalo", "Map kênh → dự án", "Webhook + Stringee"]),
    ]
    y = 110
    for name, color, steps in lanes:
        c.section_label(name, 30, y, color if isinstance(color, tuple) else color)
        for i, step in enumerate(steps):
            c.box(Box(30 + i * 390, y + 42, 360, 62, step, fill=GRAY_LIGHT if color == GRAY else GREEN_LIGHT, border=color if isinstance(color, tuple) else color))
            if i < len(steps) - 1:
                c.arrow(Arrow(390 + i * 390, y + 73, 420 + i * 390, y + 73))
        y += 130
    c.save(path)


def diagram_data_ingest(path: Path):
    c = DiagramCanvas(title="§4 — Luồng dữ liệu: Ingest → Assign → SLA → Pre-sales")
    c.header_band()
    row1 = ["Webhook\nMeta/Zalo", "Map\nform/OA?", "Tạo lead\n+ dedup", "AI score\nHot/Warm/Cold", "Auto-assign\npool NV"]
    for i, t in enumerate(row1):
        c.box(Box(40 + i * 300, 150, 260, 80, t))
        if i < 4:
            c.arrow(Arrow(300 + i * 300, 190, 340 + i * 300, 190))
    c.box(Box(980, 130, 280, 60, "Unmatched\n(chưa map)", fill=RED_LIGHT, border=ORANGE))
    c.arrow(Arrow(340, 150, 980, 160, color=ORANGE))
    c.draw.text((900, 115), "Fail", fill=ORANGE, font=c._small_font)
    row2 = ["Inbox alert\nSSE/Push", "SLA timer\nfirst-touch", "NV gọi\n+ Activity", "Gate B2\ncare ✓", "Intake\nBANT", "Pre-sales\n→ Won"]
    for i, t in enumerate(row2):
        fill = ORANGE_LIGHT if i == 1 else GREEN_LIGHT
        border = ORANGE if i == 1 else GREEN
        c.box(Box(40 + i * 250, 320, 220, 80, t, fill=fill, border=border))
        if i < 5:
            c.arrow(Arrow(260 + i * 250, 360, 290 + i * 250, 360))
    c.box(Box(40, 480, 1520, 70, "Won → Promote Customer → Service Delivery (Onboard → Deliver → Handover → Retain)", fill=BLUE_LIGHT, border=BLUE, font_size=16))
    c.save(path)


def diagram_status_gates(path: Path):
    c = DiagramCanvas(title="§5 — Trạng thái lead & Gate quan trọng")
    c.header_band()
    statuses = ["Mới", "Đã liên hệ\n(B2)", "Đang tư vấn", "Báo giá", "WON / CHỐT"]
    for i, t in enumerate(statuses):
        fill = GREEN_LIGHT if i < 4 else GREEN
        c.box(Box(60 + i * 300, 160, 240, 72, t, fill=fill))
        if i < 4:
            c.arrow(Arrow(300 + i * 300, 196, 360 + i * 300, 196))
    c.box(Box(1320, 160, 220, 72, "LOST\n(+ lý do)", fill=RED_LIGHT, border=ORANGE))
    c.arrow(Arrow(1260, 250, 1380, 232, color=ORANGE))
    c.diamond(800, 320, 70, "Deal lớn?")
    c.box(Box(560, 400, 260, 62, "Review queue\nGDKD release", fill=ORANGE_LIGHT, border=ORANGE))
    c.box(Box(920, 400, 260, 62, "Intake BANT\nGo/No-Go", fill=BLUE_LIGHT, border=BLUE))
    c.arrow(Arrow(800, 390, 690, 400))
    c.arrow(Arrow(870, 350, 1050, 400))
    gates = [
        "Gate B2 care complete",
        "Gate R5 trước báo giá",
        "GDKD approve deal lớn",
        "HĐ Active mới promote KH",
    ]
    for i, g in enumerate(gates):
        c.box(Box(60 + i * 390, 520, 350, 58, f"🔒 {g}", fill=GRAY_LIGHT, border=GRAY))
    c.save(path)


def diagram_one_page(path: Path):
    c = DiagramCanvas(title="§11 — Sơ đồ tổng hợp một trang (training)")
    c.header_band()
    c.box(Box(680, 110, 240, 58, "LEAD MỚI VÀO", fill=GREEN, border=GREEN, text_color=WHITE, font_size=17))
    c.diamond(800, 230, 65, "Loại?")
    c.box(Box(200, 320, 300, 68, "B2B Prospect\n→ Dự án PTT", fill=GREEN_LIGHT))
    c.box(Box(1100, 320, 300, 68, "CSKH spa\n→ Agency client", fill=BLUE_LIGHT, border=BLUE))
    c.arrow(Arrow(735, 168, 735, 210))
    c.arrow(Arrow(735, 295, 350, 320))
    c.arrow(Arrow(865, 295, 1250, 320))
    c.box(Box(200, 430, 300, 62, "Inbox + Gọi ≤15p", fill=GREEN_LIGHT))
    c.box(Box(1100, 430, 300, 62, "SLA board + Gọi 24h", fill=BLUE_LIGHT, border=BLUE))
    c.arrow(Arrow(350, 388, 350, 430))
    c.arrow(Arrow(1250, 388, 1250, 430))
    c.diamond(350, 540, 55, "Qualify?")
    c.diamond(1250, 540, 55, "Chốt?")
    c.box(Box(80, 630, 220, 58, "LOST", fill=RED_LIGHT, border=ORANGE))
    c.box(Box(480, 630, 280, 58, "Pre-sales → Deal Room", fill=BLUE_LIGHT, border=BLUE))
    c.box(Box(1100, 630, 300, 58, "CHỐT (spa)", fill=BLUE_LIGHT, border=BLUE))
    c.box(Box(480, 740, 280, 58, "Hub HĐ → Delivery", fill=GREEN, border=GREEN, text_color=WHITE))
    c.arrow(Arrow(350, 595, 190, 630, color=ORANGE))
    c.arrow(Arrow(410, 595, 620, 630))
    c.arrow(Arrow(1250, 595, 1250, 630, color=BLUE))
    c.arrow(Arrow(620, 688, 620, 740))
    c.save(path)


def generate_diagrams() -> dict[str, Path]:
    ASSETS.mkdir(parents=True, exist_ok=True)
    specs = {
        "01_two_flows": diagram_two_flows,
        "02_module_map": diagram_module_map,
        "03_e2e_phases": diagram_e2e_phases,
        "04_roles": diagram_roles,
        "05_data_ingest": diagram_data_ingest,
        "06_status_gates": diagram_status_gates,
        "07_one_page": diagram_one_page,
    }
    paths: dict[str, Path] = {}
    for key, fn in specs.items():
        p = ASSETS / f"diagram_{key}.png"
        fn(p)
        paths[key] = p
    return paths


# --- PowerPoint ---

PPTX_NAVY = RGBColor(0x0F, 0x17, 0x2A)
PPTX_GREEN = RGBColor(0x4A, 0x7C, 0x5C)
PPTX_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
PPTX_GRAY = RGBColor(0x64, 0x74, 0x8B)


def _blank(prs: Presentation):
    return prs.slide_layouts[6]


def _header(slide, title: str, subtitle: str = ""):
    bar = slide.shapes.add_shape(1, 0, 0, Inches(10), Inches(0.95))
    bar.fill.solid()
    bar.fill.fore_color.rgb = PPTX_NAVY
    bar.line.fill.background()
    if LOGO.is_file():
        slide.shapes.add_picture(str(LOGO), Inches(8.35), Inches(0.08), height=Inches(0.88))
    tb = slide.shapes.add_textbox(Inches(0.45), Inches(0.16), Inches(7.7), Inches(0.55))
    p = tb.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = PPTX_WHITE
    if subtitle:
        sb = slide.shapes.add_textbox(Inches(0.45), Inches(0.68), Inches(7.7), Inches(0.32))
        sp = sb.text_frame.paragraphs[0]
        sp.text = subtitle
        sp.font.size = Pt(11)
        sp.font.color.rgb = PPTX_GRAY


def _bullets(slide, items: list[str], top: float = 1.2, size: int = 16):
    box = slide.shapes.add_textbox(Inches(0.55), Inches(top), Inches(8.8), Inches(5.8))
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.font.size = Pt(size)
        p.space_after = Pt(8)


def _diagram_slide(prs: Presentation, title: str, subtitle: str, img: Path):
    slide = prs.slides.add_slide(_blank(prs))
    _header(slide, title, subtitle)
    slide.shapes.add_picture(str(img), Inches(0.35), Inches(1.05), width=Inches(9.3))


def main() -> None:
    diagrams = generate_diagrams()

    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Title
    slide = prs.slides.add_slide(_blank(prs))
    bar = slide.shapes.add_shape(1, 0, 0, Inches(10), Inches(7.5))
    bar.fill.solid()
    bar.fill.fore_color.rgb = PPTX_NAVY
    bar.line.fill.background()
    if LOGO.is_file():
        slide.shapes.add_picture(str(LOGO), Inches(3.8), Inches(0.8), height=Inches(1.2))
    tb = slide.shapes.add_textbox(Inches(0.6), Inches(2.3), Inches(8.8), Inches(2.5))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = "Leads — Sơ đồ luồng quản lý"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = PPTX_WHITE
    p2 = tf.add_paragraph()
    p2.text = "Khung bàn giao & đào tạo khách hàng"
    p2.font.size = Pt(22)
    p2.font.color.rgb = PPTX_GREEN
    p3 = tf.add_paragraph()
    p3.text = "RNOSAI · rs.pttads.vn · Phiên bản 1.0 · 2026-08-25"
    p3.font.size = Pt(14)
    p3.font.color.rgb = PPTX_GRAY

    # Intro
    slide = prs.slides.add_slide(_blank(prs))
    _header(slide, "Mục tiêu bàn giao", "Tài liệu: docs/huong-dan-su-dung/23-leads-handover-flow-and-guides.md")
    _bullets(
        slide,
        [
            "Sơ đồ luồng end-to-end từ lead mới → Won/Lost → Delivery",
            "Phân biệt 2 pipeline: B2B Sales vs CSKH vận hành (spa)",
            "Vai trò: AM · CSKH · Solution · GDKD · Marketing · IT",
            "28 file hướng dẫn chi tiết (23a–23ac) — soạn theo template bàn giao",
            "Kịch bản training 48 giờ đầu khi có lead B2B thật",
        ],
    )

    diagram_slides = [
        ("Hai luồng Lead", "Prospect B2B ≠ lead spa client đang chạy ads", "01_two_flows"),
        ("Bản đồ module Leads", "Setup → Ingest → Pipeline → Pre-sales → GDKD", "02_module_map"),
        ("Luồng chính end-to-end", "Giai đoạn 0 (Setup) → 5 (Won) + nhánh CSKH spa", "03_e2e_phases"),
        ("Luồng theo vai trò", "AM · CSKH · Solution · GDKD · Marketing/IT", "04_roles"),
        ("Luồng dữ liệu Ingest", "Webhook → map dự án → AI → assign → SLA", "05_data_ingest"),
        ("Trạng thái & Gate", "Mới → B2 → Consult → Proposal → Won + gate R5/GDKD", "06_status_gates"),
        ("Sơ đồ tổng hợp một trang", "In cho phòng training / wall poster", "07_one_page"),
    ]
    for title, subtitle, key in diagram_slides:
        _diagram_slide(prs, title, subtitle, diagrams[key])

    # 48h scenario
    slide = prs.slides.add_slide(_blank(prs))
    _header(slide, "Kịch bản 48 giờ — Lead B2B mới", "§7 — Training GDKD + AM + Marketing")
    _bullets(
        slide,
        [
            "T-7 ngày: IT map Page/Form Meta → dự án PTT · test không rơi Unmatched",
            "T+0: AM mở Inbox ≤2p → tick consent → Gọi ≤15p → log Activity",
            "T+0 ≤2h: Bắt đầu Intake BANT trên /crm/intake?lead_id=…",
            "T+1: Hoàn thành gate B2 care ✓ · GDKD release review queue (nếu có)",
            "T+1–2: Handoff Solution queue · Solution claim case",
            "T+3–5: KH MKT sơ bộ (R5) trên tab Tư vấn",
            "T+5–7: Deal Room + Proposal · GDKD approve discount nếu cần",
            "T+7–14: Hub HĐ Active → Customer 360 → Service Delivery",
        ],
        size=15,
    )

    # Feature matrix summary
    slide = prs.slides.add_slide(_blank(prs))
    _header(slide, "Ma trận tính năng (tóm tắt)", "§6 — 29 tính năng · file chi tiết 23a–23ac")
    _bullets(
        slide,
        [
            "P0 go-live: 23a Setup · 23b Dự án PTT · 23k Inbox · 23l Chi tiết · 23m Softphone",
            "P0 qualify: 23d–23f Nguồn Meta/Zalo/Khác · 23q Intake BANT",
            "P1 sales: 23r Pre-sales · 23s Solution queue · 23u Deal Room · 23p Review queue",
            "P1 gov: 23x GDKD command · 23y Speed-to-lead · 23z CSKH KPI board",
            "P2 post-won: 23v Proposal · 23w Hub HĐ · 23aa Delivery · 23ab Mobile PWA",
            "Tham chiếu đầy đủ: huong-dan-phan-he-lead-day-du.md (667 dòng)",
        ],
        size=15,
    )

    # Ops prerequisites
    slide = prs.slides.add_slide(_blank(prs))
    _header(slide, "Điều kiện vận hành", "§10 — Checklist IT trước bàn giao")
    _bullets(
        slide,
        [
            "PTT_B2B_PROJECT_OS=1 · PTT_PRESALES_ON_LEAD=1 trong .env ptt-crm-api",
            "DDL: scripts/apply_pg_ddl_b2b_lead_project_os.sh (+ wave W5 nếu Zalo thread)",
            "UAT: bash scripts/uat_b2b_project_os.sh — gate B2B-01, B2B-02",
            "RBAC: crm_leads.view/edit/assign · crm_b2b_projects.* · crm_gdkd.view_all_leads",
            "Gọi WebRTC (tuỳ chọn): PTT_B2B_CPAAS=stringee + PTT_STRINGEE_*",
            "Rollback nhanh: PTT_B2B_PROJECT_OS=0 → restart API",
        ],
        size=15,
    )

    # Closing
    slide = prs.slides.add_slide(_blank(prs))
    _header(slide, "Bước tiếp theo", "Soạn file bàn giao P0 theo template §8")
    _bullets(
        slide,
        [
            "1. Hoàn tất Setup dự án PTT + map kênh (file 23a, 23b)",
            "2. Training AM: Inbox → Gọi → Chi tiết lead (23k–23m)",
            "3. Pilot 2 deal B2B qua Deal Room (23u) trước go-live rộng",
            "4. Song song: CSKH spa theo SOP Meta 24h (23i, 23z)",
            "5. Ký nghiệm thu UAT + hypercare 2–4 tuần post go-live",
            "",
            "Liên hệ: AM dự án + Tech Lead PTT (điền tại buổi bàn giao)",
        ],
        size=16,
    )

    OUTPUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUTPUT))
    print(f"Wrote {OUTPUT}")
    print(f"Diagram PNGs: {ASSETS}/")


if __name__ == "__main__":
    main()
