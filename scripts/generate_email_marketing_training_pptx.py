#!/usr/bin/env python3
"""Generate PowerPoint: Đào tạo Email Marketing Ops trên PTTADS (ops-web)."""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
OUTPUT = ROOT / "docs" / "Email_Marketing_Ops_Training.pptx"
LOGO = ROOT / "static" / "images" / "ptt-logo.png"

NAVY = RGBColor(0x0F, 0x17, 0x2A)
GREEN = RGBColor(0x39, 0x8B, 0x43)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0x64, 0x74, 0x8B)


def _blank(prs: Presentation):
    return prs.slide_layouts[6]


def _header(slide, title: str, subtitle: str = ""):
    bar = slide.shapes.add_shape(1, 0, 0, Inches(10), Inches(0.95))
    bar.fill.solid()
    bar.fill.fore_color.rgb = NAVY
    bar.line.fill.background()
    if LOGO.is_file():
        slide.shapes.add_picture(str(LOGO), Inches(8.35), Inches(0.08), height=Inches(0.88))
    tb = slide.shapes.add_textbox(Inches(0.45), Inches(0.16), Inches(7.7), Inches(0.55))
    p = tb.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = WHITE
    if subtitle:
        sb = slide.shapes.add_textbox(Inches(0.45), Inches(0.68), Inches(7.7), Inches(0.32))
        sp = sb.text_frame.paragraphs[0]
        sp.text = subtitle
        sp.font.size = Pt(11)
        sp.font.color.rgb = GRAY


def _bullets(slide, items: list[str], top: float = 1.2):
    box = slide.shapes.add_textbox(Inches(0.55), Inches(top), Inches(8.8), Inches(5.5))
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.font.size = Pt(16)
        p.level = 0
        p.space_after = Pt(8)


def main() -> None:
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    slides: list[tuple[str, str, list[str]]] = [
        (
            "Email Marketing Ops — PTTADS",
            "ops-web /email/* · Nest API · không Flask · P1 UX parity",
            [
                "Canonical staff UI: ops-web E-01…E-13",
                "API: Nest /api/v1/email/*",
                "Workers: ptt_email/ + PostgreSQL email_mkt.*",
                "Tài liệu: docs/huong-dan-email-marketing-ops.md",
            ],
        ),
        (
            "E-01 Hub & Alerts (P1.3)",
            "Executive view + deliverability notifications",
            [
                "/email/hub — KPI sent, open rate, complaint, queue lag",
                "Alerts banner → link /email/deliverability",
                "Slack/Teams: PTT_EMAIL_SLACK_WEBHOOK / PTT_EMAIL_TEAMS_WEBHOOK",
                "Drill-down: hub → client → contacts (≤3 clicks)",
            ],
        ),
        (
            "E-07 Segment Builder (P1.2)",
            "Lifecycle · RFM · Behavior tabs",
            [
                "Tabs: Rules, Static, Lifecycle, RFM, Behavior",
                "Compute segment → member_count + eligibility filter",
                "Consent + suppression auto-excluded",
                "Route: /email/segments?client_id=",
            ],
        ),
        (
            "E-11 Deliverability & Domain Wizard (P1.5)",
            "AM-friendly onboarding",
            [
                "Wizard 3 bước: Domain → DNS records → Verify & warm-up",
                "SPF / DKIM / DMARC status badges",
                "Pause domain khi complaint cao",
                "Route: /email/deliverability",
            ],
        ),
        (
            "E-12 Reports & Grafana (P1.4)",
            "Analytics center + BI embed",
            [
                "Engagement chart, deliverability scorecard",
                "ClickHouse export job",
                "Grafana: PTT_EMAIL_GRAFANA_URL + deploy/grafana/email-ops-dashboard.json",
                "API: GET /api/v1/email/reports/bi-status",
            ],
        ),
        (
            "E-13 Governance Write (P1.1)",
            "CRUD rules + audit tail",
            [
                "Global rules: frequency_cap, quiet_hours, complaint threshold",
                "POST/PATCH/DELETE /api/v1/email/governance/rules",
                "Audit log 50 bản ghi gần nhất (before/after JSON)",
                "Cap: crm_email_mkt settings",
            ],
        ),
        (
            "Prod cutover & Gates",
            "EM-5 Gate A + P1 gate",
            [
                "Staged B1→B4: PTT_EMAIL_ENABLED → SEND → PORTAL → JOURNEYS",
                "Gate A: ./scripts/phase5_email_prod_pilot_gate.sh",
                "Handoff E2E: email-handoff.spec.ts",
                "P1 gate: ./scripts/email_p1_gate.sh",
            ],
        ),
        (
            "Checklist & tài liệu in",
            "Vận hành hàng ngày",
            [
                "Checklist A4: docs/forms/email-marketing-ops-checklist-a4.html",
                "Runbook: docs/runbooks/email-marketing-prod-pilot-checklist.md",
                "Incident: docs/runbooks/email-deliverability-incident.md",
                "Roadmap: docs/EMAIL_MARKETING_COMPLETION_ROADMAP.md",
            ],
        ),
    ]

    for title, subtitle, bullets in slides:
        slide = prs.slides.add_slide(_blank(prs))
        _header(slide, title, subtitle)
        _bullets(slide, bullets)

    OUTPUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUTPUT))
    print(f"Wrote {OUTPUT}")


if __name__ == "__main__":
    main()
