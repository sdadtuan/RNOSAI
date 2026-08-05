#!/usr/bin/env python3
"""Generate Excel + PPT: Ma trận Consult form — 12 DV + lead-gen (training AM)."""
from __future__ import annotations

import sys
from pathlib import Path

ROOT = Path(__file__).resolve().parents[1]
sys.path.insert(0, str(ROOT))

from openpyxl import Workbook
from openpyxl.styles import Alignment, Border, Font, PatternFill, Side
from openpyxl.utils import get_column_letter
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from crm_lead_intake_definitions import get_crm_field_map
from crm_svc_tasks import SERVICE_LABELS
from crm_svc_workflow_steps import SERVICE_WORKFLOW_STEPS

EXCEL_OUT = ROOT / "docs" / "exports" / "Consult_Form_Matrix_AM_Training.xlsx"
PPTX_OUT = ROOT / "docs" / "Consult_Form_Matrix_AM_Training.pptx"
LOGO = ROOT / "static" / "images" / "ptt-logo.png"

SERVICE_GROUPS: dict[str, str] = {
    "dich-vu-aeo": "Tìm kiếm tự nhiên",
    "dich-vu-seo-audit": "Tìm kiếm tự nhiên",
    "dich-vu-seo-local": "Tìm kiếm tự nhiên",
    "dich-vu-seo-tong-the": "Tìm kiếm tự nhiên",
    "quang-cao-facebook": "Quảng cáo",
    "quang-cao-google": "Quảng cáo",
    "thue-tai-khoan-quang-cao": "Quảng cáo",
    "dich-vu-quan-tri-website": "Thiết kế & web",
    "thiet-ke-website": "Thiết kế & web",
    "thiet-ke-website-tron-goi": "Thiết kế & web",
    "thiet-ke-landing-page": "Thiết kế & web",
    "tiep-thi-noi-dung": "Nội dung",
}

L2_DOCS: dict[str, str] = {
    "dich-vu-aeo": "URL, content hiện có, test query brand trên ChatGPT/Gemini/Perplexity",
    "dich-vu-seo-audit": "GSC read, GA4, hosting, mục tiêu audit",
    "dich-vu-seo-local": "Link GBP, NAP chi nhánh, ảnh cửa hàng, review count",
    "dich-vu-seo-tong-the": "GSC, GA4, 2–3 đối thủ, danh sách từ khóa seed",
    "quang-cao-facebook": "Ads account read, pixel, LP URL, lịch sử spend",
    "quang-cao-google": "Account read, conversion tracking, LP, CPC ước tính",
    "thue-tai-khoan-quang-cao": "Lịch sử policy, sản phẩm QC, landing compliance",
    "dich-vu-quan-tri-website": "Admin WP/hosting, backup status, plugin list",
    "thiet-ke-website": "Brand assets, sitemap draft, website tham khảo (URLs)",
    "thiet-ke-website-tron-goi": "Feature list, payment/CRM, hosting/domain",
    "thiet-ke-landing-page": "Offer, copy draft, campaign đi kèm, brand guideline",
    "tiep-thi-noi-dung": "Content hiện có, brand voice, competitor URLs",
}

CHECKLIST_BY_GROUP: dict[str, str] = {
    "Tìm kiếm tự nhiên": "Scope audit rõ · baseline đo được · đối thủ/KW đã liệt kê",
    "Quảng cáo": "Objective + targeting draft · KPI target (CPL/ROAS) · kế hoạch account access",
    "Thiết kế & web": "Deliverables (design/code/pages) · số revision · timeline",
    "Nội dung": "Số bài/tháng · kênh · tone of voice · calendar tháng 1 sơ bộ",
    "Pilot / inbound": "Funnel end-to-end · KPI cam kết draft · scope PTT vs KH tự làm",
}

LEAD_GEN_SLUG = "lead-gen"
LEAD_GEN_GROUP = "Performance / inbound"
LEAD_GEN_LABEL = "Lead generation"
LEAD_GEN_TASK = "Discovery funnel & kênh lead generation"
LEAD_GEN_FIELDS: list[dict[str, str]] = [
    {"key": "current_status", "label": "Hiện trạng funnel & kênh", "type": "textarea"},
    {"key": "target_audience", "label": "ICP & đối tượng mục tiêu", "type": "textarea"},
    {"key": "conversion_metrics", "label": "KPI hiện tại & mục tiêu", "type": "textarea"},
    {"key": "scope_recommendation", "label": "Phạm vi PTT đề xuất", "type": "textarea"},
]
LEAD_GEN_L2 = "Meta lead form export, Ads account read, LP URL, CRM screenshot, spend 3 tháng"
LEAD_GEN_PREFILL = (
    "need→current_status · niche→target_audience · "
    "budget/daily_budget/monthly_budget→conversion_metrics · "
    "campaign_goal→scope_recommendation · discovery_responses→current_status (Phase 2)"
)

# Excel styling
HEADER_FILL = PatternFill("solid", fgColor="4F46E5")
HEADER_FONT = Font(bold=True, color="FFFFFF", size=11)
TITLE_FONT = Font(bold=True, size=14, color="4F46E5")
SUB_FONT = Font(italic=True, size=10, color="555555")
WARN_FILL = PatternFill("solid", fgColor="FEF3C7")
THIN = Side(style="thin", color="CCCCCC")
BORDER = Border(left=THIN, right=THIN, top=THIN, bottom=THIN)
WRAP = Alignment(wrap_text=True, vertical="top")

# PPT styling
NAVY = RGBColor(0x0F, 0x17, 0x2A)
ACCENT = RGBColor(0x63, 0x66, 0xF1)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0x64, 0x74, 0x8B)
DARK_TEXT = RGBColor(0x1E, 0x29, 0x3B)
LIGHT_BG = RGBColor(0xF1, 0xF5, 0xF9)
ACCENT_LIGHT = RGBColor(0xC7, 0xD2, 0xFE)


def _sorted_slugs() -> list[str]:
    return sorted(SERVICE_WORKFLOW_STEPS.keys())


def _consult_step(slug: str) -> dict | None:
    steps = SERVICE_WORKFLOW_STEPS.get(slug, {}).get("consult") or []
    return steps[0] if steps else None


def _fields_str(slug: str) -> str:
    step = _consult_step(slug)
    if not step:
        return "—"
    parts = [f"{f['key']} ({f['label']})" for f in (step.get("form_fields") or [])]
    return " · ".join(parts)


def _prefill_str(slug: str) -> str:
    m = get_crm_field_map(slug)
    if not m:
        return "—"
    return " · ".join(f"{k}→{v}" for k, v in sorted(m.items()))


def _field_keys_short(slug: str) -> str:
    if slug == LEAD_GEN_SLUG:
        return " · ".join(f["key"] for f in LEAD_GEN_FIELDS)
    step = _consult_step(slug)
    if not step:
        return "—"
    return ", ".join(f["key"] for f in (step.get("form_fields") or []))


def _matrix_rows(include_lead_gen: bool = True) -> list[list[str]]:
    rows: list[list[str]] = []
    for i, slug in enumerate(_sorted_slugs(), 1):
        step = _consult_step(slug)
        if not step:
            continue
        rows.append([
            str(i),
            SERVICE_GROUPS.get(slug, "—"),
            SERVICE_LABELS.get(slug, slug),
            slug,
            step.get("title", ""),
            _fields_str(slug),
            L2_DOCS.get(slug, "—"),
            _prefill_str(slug),
        ])
    if include_lead_gen:
        rows.append([
            str(len(rows) + 1),
            "Performance / inbound",
            LEAD_GEN_LABEL,
            LEAD_GEN_SLUG,
            LEAD_GEN_TASK,
            " · ".join(f"{f['key']} ({f['label']})" for f in LEAD_GEN_FIELDS),
            LEAD_GEN_L2,
            LEAD_GEN_PREFILL,
        ])
    return rows


# ── Excel ──────────────────────────────────────────────────────────────────

def _style_header_row(ws, row: int, cols: int) -> None:
    for c in range(1, cols + 1):
        cell = ws.cell(row=row, column=c)
        cell.fill = HEADER_FILL
        cell.font = HEADER_FONT
        cell.alignment = Alignment(horizontal="center", vertical="center", wrap_text=True)
        cell.border = BORDER


def _write_table(ws, start_row: int, headers: list[str], rows: list[list[str]]) -> int:
    for i, h in enumerate(headers, 1):
        ws.cell(row=start_row, column=i, value=h)
    _style_header_row(ws, start_row, len(headers))
    r = start_row + 1
    for row in rows:
        for i, val in enumerate(row, 1):
            cell = ws.cell(row=r, column=i, value=val)
            cell.alignment = WRAP
            cell.border = BORDER
            if LEAD_GEN_SLUG in str(row[3] if len(row) > 3 else ""):
                cell.fill = WARN_FILL
        r += 1
    return r


def _autosize(ws, max_width: int = 52) -> None:
    for col in ws.columns:
        letter = get_column_letter(col[0].column)
        width = max(len(str(c.value or "")) for c in col)
        ws.column_dimensions[letter].width = min(max(width + 2, 10), max_width)


def build_excel() -> Path:
    EXCEL_OUT.parent.mkdir(parents=True, exist_ok=True)
    wb = Workbook()

    # Sheet 1 — Tổng quan
    ws = wb.active
    ws.title = "1-Tong_quan"
    ws["A1"] = "MA TRẬN CONSULT FORM — TRAINING AM"
    ws["A1"].font = TITLE_FONT
    ws["A2"] = (
        "Phiên bản 1.0 · 2026-08-05 · 12 dịch vụ CRM + lead-gen (pilot) · "
        "Nguồn: crm_svc_workflow_steps.py + get_crm_field_map()"
    )
    ws["A2"].font = SUB_FONT
    overview = [
        ["Kiến trúc", "3 lớp", "A: Core audit · B: Form theo DV · C: Gate R5 KH MKT sơ bộ"],
        ["AI prompt", "consult_analysis", "Context: BANT, decision, Intake recap, red flags"],
        ["Exit Consult", "Task ✓ + R5", "CTA Chuyển → Báo giá trên funnel stepper (G4)"],
        ["SLA", "≤48h", "Sau meeting Consult → chuyển Proposal"],
        ["SOP", "consult-stage-am-sop.md", "docs/runbooks/"],
        ["Tạo lại file", "python3 scripts/generate_consult_form_matrix_training.py", ""],
    ]
    _write_table(ws, 4, ["Khái niệm", "Giá trị", "Mô tả"], overview)
    _autosize(ws, 40)

    # Sheet 2 — Ma trận 13 DV
    ws2 = wb.create_sheet("2-Ma_tran_13_DV")
    ws2["A1"] = "MA TRẬN TỔNG — CONSULT FORM THEO DỊCH VỤ"
    ws2["A1"].font = TITLE_FONT
    headers = ["#", "Nhóm", "Dịch vụ", "Slug", "Task Consult", "Form fields", "Tài liệu L2", "Prefill Intake→Consult"]
    _write_table(ws2, 3, headers, _matrix_rows())
    _autosize(ws2)

    # Sheet 3 — Form fields chi tiết
    ws3 = wb.create_sheet("3-Form_fields")
    ws3["A1"] = "CHI TIẾT FORM FIELDS — TASK CONSULT"
    ws3["A1"].font = TITLE_FONT
    detail_rows: list[list[str]] = []
    for slug in _sorted_slugs():
        step = _consult_step(slug)
        if not step:
            continue
        for f in step.get("form_fields") or []:
            detail_rows.append([
                SERVICE_GROUPS.get(slug, ""),
                SERVICE_LABELS.get(slug, slug),
                slug,
                step.get("title", ""),
                f.get("key", ""),
                f.get("label", ""),
                f.get("type", ""),
                "Bắt buộc trước ✓",
            ])
    for f in LEAD_GEN_FIELDS:
        detail_rows.append([
            "Performance / inbound",
            LEAD_GEN_LABEL,
            LEAD_GEN_SLUG,
            LEAD_GEN_TASK,
            f["key"],
            f["label"],
            f["type"],
            "Shipped CRM template",
        ])
    _write_table(
        ws3,
        3,
        ["Nhóm", "Dịch vụ", "Slug", "Task", "Field key", "Label", "Loại", "Ghi chú"],
        detail_rows,
    )
    _autosize(ws3)

    # Sheet 4 — Prefill
    ws4 = wb.create_sheet("4-Prefill_Intake")
    ws4["A1"] = "PREFILL INTAKE / LEAD → CONSULT (C2)"
    ws4["A1"].font = TITLE_FONT
    prefill_rows: list[list[str]] = []
    for slug in _sorted_slugs():
        m = get_crm_field_map(slug)
        for src, tgt in sorted(m.items()):
            prefill_rows.append([
                SERVICE_LABELS.get(slug, slug),
                slug,
                src,
                tgt,
                "Auto khi Chuyển → Tư vấn hoặc nút Prefill",
            ])
    for part in LEAD_GEN_PREFILL.replace(" ⚠️ đề xuất", "").split(" · "):
        if "→" in part:
            src, tgt = part.split("→", 1)
            prefill_rows.append([LEAD_GEN_LABEL, LEAD_GEN_SLUG, src.strip(), tgt.strip(), "Auto prefill"])
    _write_table(ws4, 3, ["Dịch vụ", "Slug", "Intake/Lead field", "Consult field", "Ghi chú"], prefill_rows)
    _autosize(ws4)

    # Sheet 5 — Gate R5
    ws5 = wb.create_sheet("5-Gate_R5_KH_MKT")
    ws5["A1"] = "GATE G4 — KH MARKETING SƠ BỘ (R5)"
    ws5["A1"].font = TITLE_FONT
    r5_rows = [
        ["name", "Tên kế hoạch MKT sơ bộ", "Bắt buộc", "Lead panel → Lưu KH MKT sơ bộ"],
        ["north_star", "North Star", "north_star HOẶC objectives", ""],
        ["objectives", "Mục tiêu chiến lược", "north_star HOẶC objectives", ""],
        ["market_message", "Thông điệp thị trường (strategy_framework)", "Bắt buộc", ""],
        ["media_reach", "Kênh tiếp cận", "Bắt buộc", ""],
        ["conversion_strategy", "Chiến lược chuyển đổi", "Bắt buộc", ""],
    ]
    _write_table(ws5, 3, ["Field key", "Label", "Rule", "Thao tác CRM"], r5_rows)
    _autosize(ws5, 36)

    # Sheet 6 — Checklist theo nhóm
    ws6 = wb.create_sheet("6-Checklist_nhom")
    ws6["A1"] = "CHECKLIST ✓ CONSULT → PROPOSAL THEO NHÓM"
    ws6["A1"].font = TITLE_FONT
    chk_rows = [[g, c, "Task Consult ✓ + R5 đủ → Chuyển → Báo giá"] for g, c in CHECKLIST_BY_GROUP.items()]
    _write_table(ws6, 3, ["Nhóm dịch vụ", "AM xác nhận trước ✓", "Gate"], chk_rows)
    _autosize(ws6, 44)

    # Sheet 7 — Quy trình 5 bước
    ws7 = wb.create_sheet("7-Quy_trinh_5_buoc")
    ws7["A1"] = "QUY TRÌNH CONSULT — 5 BƯỚC AM"
    ws7["A1"].font = TITLE_FONT
    flow_rows = [
        ["1", "Đọc brief", "Consult Brief + Intake completed", "Ngay khi vào Consult"],
        ["2", "Thu tài liệu L2", "Theo sheet 2 — cột Tài liệu L2", "1–3 ngày trước meeting"],
        ["3", "Audit/discovery", "Điền form task Consult trên Lead #funnel-presales", "3–7 ngày"],
        ["4", "AI Hỗ trợ", "consult_analysis → review notes", "Cùng buổi tư vấn"],
        ["5", "✓ + R5 + Proposal", "Tick task · KH MKT sơ bộ · Chuyển → Báo giá", "≤48h sau meeting"],
    ]
    _write_table(ws7, 3, ["Bước", "Việc làm", "Chi tiết", "SLA"], flow_rows)
    _autosize(ws7, 36)

    wb.save(EXCEL_OUT)
    return EXCEL_OUT


# ── PPT ────────────────────────────────────────────────────────────────────

def _logo_exists() -> bool:
    return LOGO.is_file()


def _blank(prs: Presentation):
    return prs.slide_layouts[6]


def _header(slide, title: str, subtitle: str = ""):
    bar = slide.shapes.add_shape(1, 0, 0, Inches(10), Inches(0.95))
    bar.fill.solid()
    bar.fill.fore_color.rgb = NAVY
    bar.line.fill.background()
    if _logo_exists():
        slide.shapes.add_picture(str(LOGO), Inches(8.35), Inches(0.08), height=Inches(0.88))
    tb = slide.shapes.add_textbox(Inches(0.45), Inches(0.16), Inches(7.7), Inches(0.55))
    p = tb.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = WHITE
    if subtitle:
        sb = slide.shapes.add_textbox(Inches(0.45), Inches(0.68), Inches(7.7), Inches(0.32))
        sp = sb.text_frame.paragraphs[0]
        sp.text = subtitle
        sp.font.size = Pt(11)
        sp.font.color.rgb = ACCENT_LIGHT


def _bullets(slide, items: list[str], top=1.15, size=13):
    box = slide.shapes.add_textbox(Inches(0.45), Inches(top), Inches(9.1), Inches(6.0))
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.font.size = Pt(size)
        p.font.color.rgb = DARK_TEXT
        p.space_after = Pt(4)


def _table_slide(
    prs,
    title: str,
    headers: list[str],
    rows: list[list[str]],
    subtitle: str = "",
    col_widths: list[float] | None = None,
):
    slide = prs.slides.add_slide(_blank(prs))
    _header(slide, title, subtitle)
    nr, nc = len(rows) + 1, len(headers)
    tbl_shape = slide.shapes.add_table(nr, nc, Inches(0.2), Inches(1.1), Inches(9.6), Inches(5.9))
    tbl = tbl_shape.table
    if col_widths and len(col_widths) == nc:
        total = sum(col_widths)
        for j, w in enumerate(col_widths):
            tbl.columns[j].width = Inches(9.6 * w / total)
    for j, h in enumerate(headers):
        c = tbl.cell(0, j)
        c.text = h
        for p in c.text_frame.paragraphs:
            p.font.bold = True
            p.font.size = Pt(9)
            p.font.color.rgb = WHITE
        c.fill.solid()
        c.fill.fore_color.rgb = ACCENT
    for i, row in enumerate(rows, 1):
        for j, val in enumerate(row):
            c = tbl.cell(i, j)
            c.text = val
            for p in c.text_frame.paragraphs:
                p.font.size = Pt(8)
                p.font.color.rgb = DARK_TEXT
            if i % 2 == 0:
                c.fill.solid()
                c.fill.fore_color.rgb = LIGHT_BG


def build_pptx() -> Path:
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Title
    slide = prs.slides.add_slide(_blank(prs))
    bg = slide.shapes.add_shape(1, 0, 0, Inches(10), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = NAVY
    bg.line.fill.background()
    if _logo_exists():
        slide.shapes.add_picture(str(LOGO), Inches(4.1), Inches(0.85), height=Inches(1.35))
    accent = slide.shapes.add_shape(1, 0, Inches(2.35), Inches(10), Inches(0.06))
    accent.fill.solid()
    accent.fill.fore_color.rgb = ACCENT
    accent.line.fill.background()
    for y, text, sz, bold in [
        (2.65, "Ma trận Consult Form", 32, True),
        (3.25, "Training AM — 12 dịch vụ + lead-gen", 18, False),
        (3.85, "Form · Prefill · Tài liệu L2 · Gate R5 · Checklist", 14, False),
        (4.45, "2026-08-05 · rs.pttads.vn · Funnel stepper Phase 2.5", 11, False),
    ]:
        tb = slide.shapes.add_textbox(Inches(0.6), Inches(y), Inches(8.8), Inches(0.55))
        p = tb.text_frame.paragraphs[0]
        p.text = text
        p.font.size = Pt(sz)
        p.font.bold = bold
        p.font.color.rgb = WHITE if sz >= 18 else GRAY
        p.alignment = PP_ALIGN.CENTER

    slide = prs.slides.add_slide(_blank(prs))
    _header(slide, "Mục lục")
    _bullets(slide, [
        "I. Kiến trúc 3 lớp Consult form",
        "II. Quy trình 5 bước AM",
        "III. Ma trận 13 dịch vụ (form + L2)",
        "IV. lead-gen (Meta inbound default slug)",
        "V. Prefill Intake → Consult",
        "VI. Gate R5 — KH MKT sơ bộ",
        "VII. Checklist theo nhóm · SLA",
        "VIII. Tài liệu repo & tạo lại file",
    ], size=14)

    slide = prs.slides.add_slide(_blank(prs))
    _header(slide, "Kiến trúc 3 lớp")
    _bullets(slide, [
        "Lớp A — Core: ghi nhận audit/discovery (mọi DV)",
        "Lớp B — Theo dịch vụ: 2–4 field task Consult (sheet Form fields)",
        "Lớp C — Gate R5: KH Marketing sơ bộ trước Chuyển → Báo giá (G4)",
        "",
        "AI: consult_analysis — context BANT, decision, Intake, red flags",
        "Exit: Task Consult ✓ + R5 đủ → CTA Chuyển → Báo giá trên stepper",
        "",
        "Lead #900000002 (pilot): service_slug = lead-gen — cần template riêng",
    ], size=12)

    slide = prs.slides.add_slide(_blank(prs))
    _header(slide, "Quy trình Consult — 5 bước")
    _bullets(slide, [
        "1. Đọc Consult Brief + Intake completed — ngay khi vào Consult",
        "2. Thu tài liệu L2 theo dịch vụ — 1–3 ngày trước meeting",
        "3. Buổi audit/discovery — điền form task Consult trên Lead",
        "4. AI Hỗ trợ (consult_analysis) — review trong buổi",
        "5. Tick ✓ Consult + KH MKT sơ bộ → Chuyển → Báo giá (≤48h)",
        "",
        "Handoff Intake: stepper CTA Chuyển → Tư vấn (Phase 2.5)",
        "Không nhập lại niche/budget — dùng Prefill hoặc đọc Brief",
    ], size=12)

    matrix = _matrix_rows()
    mid = (len(matrix) + 1) // 2
    compact_headers = ["#", "Nhóm", "Dịch vụ", "Task Consult", "Form (keys)"]
    for chunk, title in [(matrix[:mid], "Ma trận Consult — phần 1"), (matrix[mid:], "Ma trận Consult — phần 2")]:
        rows = [[r[0], r[1], r[2], r[4], _field_keys_short(r[3])] for r in chunk]
        _table_slide(
            prs,
            title,
            compact_headers,
            rows,
            subtitle="Chi tiết đầy đủ: Excel sheet 2-Ma_tran_13_DV",
            col_widths=[0.35, 1.0, 1.3, 2.2, 2.5],
        )

    slide = prs.slides.add_slide(_blank(prs))
    _header(slide, "lead-gen — Template Meta inbound", "Slug catalog default · crm_svc_workflow_steps.py")
    _bullets(slide, [
        f"Task: {LEAD_GEN_TASK}",
        "",
        "Form fields:",
        "  current_status — Hiện trạng funnel & kênh",
        "  target_audience — ICP & đối tượng mục tiêu",
        "  conversion_metrics — KPI hiện tại & mục tiêu",
        "  scope_recommendation — Phạm vi PTT đề xuất",
        "",
        f"L2: {LEAD_GEN_L2}",
        f"Prefill: {LEAD_GEN_PREFILL}",
    ], size=11)

    prefill_sample = []
    for slug in ("dich-vu-seo-tong-the", "quang-cao-facebook", "quang-cao-google", LEAD_GEN_SLUG):
        if slug == LEAD_GEN_SLUG:
            prefill_sample.append([LEAD_GEN_LABEL, LEAD_GEN_PREFILL[:90]])
        else:
            prefill_sample.append([SERVICE_LABELS.get(slug, slug), _prefill_str(slug)[:90]])
    _table_slide(
        prs,
        "Prefill Intake → Consult (mẫu)",
        ["Dịch vụ", "Mapping (Lead/Intake → Consult field)"],
        prefill_sample,
        subtitle="Auto khi Chuyển → Tư vấn · get_crm_field_map(slug)",
        col_widths=[1.2, 3.5],
    )

    _table_slide(
        prs,
        "Gate G4 — KH Marketing sơ bộ (R5)",
        ["Field", "Rule"],
        [
            ["name", "Bắt buộc"],
            ["north_star hoặc objectives", "Ít nhất một"],
            ["market_message", "Bắt buộc (strategy_framework)"],
            ["media_reach", "Bắt buộc"],
            ["conversion_strategy", "Bắt buộc"],
        ],
        subtitle="Lead panel → Lưu KH MKT sơ bộ → CTA Chuyển → Báo giá",
        col_widths=[1.5, 2.5],
    )

    _table_slide(
        prs,
        "Checklist ✓ Consult → Proposal",
        ["Nhóm", "AM xác nhận"],
        [[g, c] for g, c in CHECKLIST_BY_GROUP.items()],
        col_widths=[1.0, 3.5],
    )

    slide = prs.slides.add_slide(_blank(prs))
    _header(slide, "Tài liệu & tạo lại")
    _bullets(slide, [
        "Excel: docs/exports/Consult_Form_Matrix_AM_Training.xlsx",
        "PPT: docs/Consult_Form_Matrix_AM_Training.pptx",
        "SOP: docs/runbooks/consult-stage-am-sop.md",
        "Runbook task: docs/runbooks/consult-stage-service-tasks.md",
        "Training 45p: docs/runbooks/consult-stage-training-guide.md",
        "",
        "Tạo lại:",
        "  python3 scripts/generate_consult_form_matrix_training.py",
        "",
        "Pilot AM: in Excel sheet 6-Checklist · tick khi training",
    ], size=12)

    # Closing
    slide = prs.slides.add_slide(_blank(prs))
    bg = slide.shapes.add_shape(1, 0, 0, Inches(10), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = NAVY
    bg.line.fill.background()
    tb = slide.shapes.add_textbox(Inches(1), Inches(3.0), Inches(8), Inches(1.2))
    p = tb.text_frame.paragraphs[0]
    p.text = "Cảm ơn · PTT Advertising Solutions"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    PPTX_OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(PPTX_OUT))
    return PPTX_OUT


def main() -> None:
    xlsx = build_excel()
    pptx = build_pptx()
    print(f"Wrote {xlsx}")
    print(f"Wrote {pptx}")


if __name__ == "__main__":
    main()
