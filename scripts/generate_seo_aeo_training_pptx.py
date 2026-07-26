#!/usr/bin/env python3
"""Generate PowerPoint: Đào tạo SEO/AEO Ops trên PTTADS (ops-web)."""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
OUTPUT = ROOT / "docs" / "SEO_AEO_Ops_Training.pptx"
LOGO = ROOT / "static" / "images" / "ptt-logo.png"

NAVY = RGBColor(0x0F, 0x17, 0x2A)
GREEN = RGBColor(0x39, 0x8B, 0x43)
GREEN_LIGHT = RGBColor(0xDC, 0xFC, 0xE7)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0x64, 0x74, 0x8B)
DARK = RGBColor(0x1E, 0x29, 0x3B)
LIGHT = RGBColor(0xF1, 0xF5, 0xF9)


def _blank(prs: Presentation):
    return prs.slide_layouts[6]


def _logo_ok() -> bool:
    return LOGO.is_file()


def _header(slide, title: str, subtitle: str = ""):
    bar = slide.shapes.add_shape(1, 0, 0, Inches(10), Inches(0.95))
    bar.fill.solid()
    bar.fill.fore_color.rgb = NAVY
    bar.line.fill.background()
    if _logo_ok():
        slide.shapes.add_picture(str(LOGO), Inches(8.35), Inches(0.08), height=Inches(0.88))
    tb = slide.shapes.add_textbox(Inches(0.45), Inches(0.16), Inches(7.7), Inches(0.55))
    p = tb.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = WHITE
    if subtitle:
        sb = slide.shapes.add_textbox(Inches(0.45), Inches(0.68), Inches(7.7), Inches(0.32))
        sp = sb.text_frame.paragraphs[0]
        sp.text = subtitle
        sp.font.size = Pt(11)
        sp.font.color.rgb = GREEN_LIGHT


def _bullets(slide, items: list[str], top=1.15, size=13):
    box = slide.shapes.add_textbox(Inches(0.45), Inches(top), Inches(9.1), Inches(6.0))
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.level = 1 if item.startswith("  ") else 0
        p.font.size = Pt(size - 1 if p.level else size)
        p.font.color.rgb = GRAY if p.level else DARK
        p.space_after = Pt(4)


def _table_slide(
    prs,
    title: str,
    headers: list[str],
    rows: list[list[str]],
    subtitle: str = "",
    col_widths: list[float] | None = None,
):
    slide = prs.slides.add_slide(_blank(prs))
    _header(slide, title, subtitle)
    tbl_shape = slide.shapes.add_table(
        len(rows) + 1, len(headers), Inches(0.25), Inches(1.12), Inches(9.5), Inches(5.85)
    )
    tbl = tbl_shape.table
    if col_widths and len(col_widths) == len(headers):
        total = sum(col_widths)
        for j, w in enumerate(col_widths):
            tbl.columns[j].width = Inches(9.5 * w / total)
    for j, h in enumerate(headers):
        c = tbl.cell(0, j)
        c.text = h
        for p in c.text_frame.paragraphs:
            p.font.bold = True
            p.font.size = Pt(10)
            p.font.color.rgb = WHITE
        c.fill.solid()
        c.fill.fore_color.rgb = GREEN
    for i, row in enumerate(rows, 1):
        for j, val in enumerate(row):
            c = tbl.cell(i, j)
            c.text = val
            for p in c.text_frame.paragraphs:
                p.font.size = Pt(9)
                p.font.color.rgb = DARK
            if i % 2 == 0:
                c.fill.solid()
                c.fill.fore_color.rgb = LIGHT


def _title_slide(prs: Presentation):
    slide = prs.slides.add_slide(_blank(prs))
    bg = slide.shapes.add_shape(1, 0, 0, Inches(10), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = NAVY
    bg.line.fill.background()
    if _logo_ok():
        slide.shapes.add_picture(str(LOGO), Inches(4.1), Inches(0.85), height=Inches(1.35))
    accent = slide.shapes.add_shape(1, 0, Inches(2.35), Inches(10), Inches(0.06))
    accent.fill.solid()
    accent.fill.fore_color.rgb = GREEN
    accent.line.fill.background()
    lines = [
        (2.65, "SEO/AEO Enterprise Ops", 34, True, WHITE),
        (3.35, "Đào tạo vận hành trên PTTADS · ops-web", 18, False, GREEN_LIGHT),
        (4.05, "Hub · Research · Content · Technical · AEO · Báo cáo", 14, False, GRAY),
        (4.65, "v2.0 · 2026-07-25 · ~45 phút · docs/huong-dan-seo-aeo-ops.md", 11, False, GRAY),
    ]
    for y, text, sz, bold, color in lines:
        tb = slide.shapes.add_textbox(Inches(0.6), Inches(y), Inches(8.8), Inches(0.55))
        p = tb.text_frame.paragraphs[0]
        p.text = text
        p.font.size = Pt(sz)
        p.font.bold = bold
        p.font.color.rgb = color
        p.alignment = PP_ALIGN.CENTER


def _closing_slide(prs: Presentation):
    slide = prs.slides.add_slide(_blank(prs))
    bg = slide.shapes.add_shape(1, 0, 0, Inches(10), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = NAVY
    bg.line.fill.background()
    if _logo_ok():
        slide.shapes.add_picture(str(LOGO), Inches(4.25), Inches(2.0), height=Inches(1.2))
    for y, text, sz in [
        (3.45, "Cảm ơn", 40),
        (4.35, "Checklist A4: docs/forms/seo-aeo-ops-checklist-a4.html", 13),
        (4.85, "PTT Advertising Solutions · SEO/AEO Ops Training", 12),
    ]:
        tb = slide.shapes.add_textbox(Inches(0.8), Inches(y), Inches(8.4), Inches(0.6))
        p = tb.text_frame.paragraphs[0]
        p.text = text
        p.font.size = Pt(sz)
        p.font.bold = sz >= 20
        p.font.color.rgb = WHITE if sz >= 20 else GRAY
        p.alignment = PP_ALIGN.CENTER


def build() -> Path:
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    _title_slide(prs)

    slide = prs.slides.add_slide(_blank(prs))
    _header(slide, "Mục lục")
    _bullets(
        slide,
        [
            "I. Tổng quan phân hệ & stack kỹ thuật",
            "II. Đăng nhập, phân quyền, điều hướng",
            "III. Hub & Client workspace (S-01 → S-04)",
            "IV. Chiến lược OKR & Research (S-05 → S-06)",
            "V. Content pipeline & Publish (S-07 → S-08)",
            "VI. Technical · AEO · Báo cáo (S-09 → S-12)",
            "VII. Governance · Automations · Rank · Freshness",
            "VIII. Luồng hàng ngày & onboard client",
            "IX. Go-live Gate A · Tài liệu tham khảo",
        ],
        top=1.15,
        size=14,
    )

    slide = prs.slides.add_slide(_blank(prs))
    _header(slide, "Vision — Vòng đời SEO/AEO", "SPEC_SEO_AEO_OPERATING_SYSTEM.md §1")
    _bullets(
        slide,
        [
            "Chiến lược → Nghiên cứu → Sản xuất nội dung → QA kỹ thuật",
            "→ Tối ưu AEO → Publish → Giám sát → Refresh → Báo cáo",
            "",
            "Mục tiêu agency:",
            "  • Tăng organic traffic & non-brand visibility",
            "  • Tăng citations trong AI / answer engines",
            "  • Giảm thời gian vận hành · tăng QA pass rate",
            "",
            "Nguyên tắc: multi-tenant · workflow-driven · measurable · governance-heavy",
        ],
    )

    slide = prs.slides.add_slide(_blank(prs))
    _header(slide, "Stack canonical (2026-07-25)", "Không dùng Flask /crm/seo cho staff")
    _bullets(
        slide,
        [
            "Staff UI: ops-web https://rs.pttads.vn/seo/*",
            "Staff API: Nest ptt-crm-api /api/v1/seo/*",
            "Domain/workers: Python ptt_seo/ + ptt_worker",
            "Data: PostgreSQL schema seo_aeo.* (PostgreSQL-only)",
            "CRM master: SQLite crm_customers (customer_id)",
            "",
            "Flask /crm/seo/* → nginx redirect sang /seo/*",
            "Portal client: portal.pttads.vn/seo (flag PTT_PORTAL_SEO_ENABLED)",
        ],
    )

    _table_slide(
        prs,
        "Bản đồ module ops-web",
        ["Screen", "Route", "Vai trò chính"],
        [
            ["S-01 Hub", "/seo/hub", "KPI tổng · client health · drill-down"],
            ["S-02/03 Clients", "/seo/clients/:id", "Workspace · settings · nav module"],
            ["S-05 Strategy", "/seo/strategy", "OKR/KPI tree · refresh · editor"],
            ["S-06 Research", "/seo/research", "7 tabs · brief → content"],
            ["S-07 Content", "/seo/content", "Kanban 13 stage · filters"],
            ["S-09 Technical", "/seo/technical", "GSC/GA4 OAuth · crawl · CWV"],
            ["S-10 AEO", "/seo/aeo", "Question bank · batch scan"],
            ["S-12 Reports", "/seo/reports", "Dashboard · attribution · PDF"],
            ["S-14 Governance", "/seo/governance", "Policy publish"],
            ["Gate A", "/seo/gate-a", "Go-live readiness"],
        ],
        subtitle="19 routes · chi tiết: huong-dan-seo-aeo-ops.md §6",
        col_widths=[1.0, 1.6, 3.0],
    )

    _table_slide(
        prs,
        "Personas & phân quyền",
        ["Vai trò", "Section keys", "Không được"],
        [
            ["Head SEO / MKT-01", "Cả 6 keys crm_seo_aeo_*", "—"],
            ["Strategist", "view + write + reports", "approve"],
            ["Writer", "view + write", "approve · settings"],
            ["Tech SEO", "view + technical", "approve content"],
            ["AM / KD-01", "view + settings + reports", "technical import"],
        ],
        subtitle="Admin → Phân quyền trang · ops-web ẩn nút nếu thiếu cap",
        col_widths=[1.4, 2.2, 1.8],
    )

    slide = prs.slides.add_slide(_blank(prs))
    _header(slide, "Đăng nhập & điều hướng")
    _bullets(
        slide,
        [
            "1. https://rs.pttads.vn/login (staff ops-web)",
            "2. Sidebar → Agency & Hub → SEO/AEO Hub",
            "3. Bookmark cũ /crm/seo tự redirect → /seo/hub",
            "",
            "Context bar client workspace:",
            "  Client name · Domain · Market · Owner · Tier",
            "",
            "Badge cảnh báo (UI/UX spec §3.4):",
            "  • Critical issues → tab Kỹ thuật đỏ",
            "  • Content overdue → tab Nội dung cam",
            "  • AEO coverage <50% → tab AEO vàng",
            "  • Sync failed → banner Hub",
        ],
    )

    slide = prs.slides.add_slide(_blank(prs))
    _header(slide, "S-01 Executive Hub", "/seo/hub")
    _bullets(
        slide,
        [
            "KPI grid: clients · GSC clicks · content delivery · critical · AEO",
            "Bảng client health → click vào workspace (≤3 click)",
            "Drill-down P2: content → /seo/content · technical → /seo/technical",
            "",
            "Thao tác hàng ngày (Head SEO):",
            "  1. Kiểm tra sync banner (GSC/GA4 OK?)",
            "  2. Review client health đỏ/vàng",
            "  3. Drill-down issue critical",
            "  4. Export / báo cáo executive nếu cần",
        ],
    )

    slide = prs.slides.add_slide(_blank(prs))
    _header(slide, "S-03 Client workspace", "/seo/clients/:id · SeoClientWorkspaceNav")
    _bullets(
        slide,
        [
            "3 tab nội bộ: Tổng quan · Roadmap · Tasks",
            "7 link module với ?customer_id=:",
            "  Nghiên cứu · Nội dung · Kỹ thuật · AEO · Authority · Báo cáo",
            "",
            "S-04 Settings (cap crm_seo_aeo_settings):",
            "  • Domains / markets / languages",
            "  • Brand guidelines JSON",
            "  • GSC site + GA4 property (OAuth)",
            "  • CMS webhook URL + secret",
            "  • Approvers client_review chain",
        ],
    )

    slide = prs.slides.add_slide(_blank(prs))
    _header(slide, "S-05 Chiến lược OKR", "/seo/strategy")
    _bullets(
        slide,
        [
            "Cây Goal → KPI → Initiative",
            "Refresh KPI: cập nhật current_value từ GSC/GA4 live",
            "P1: form tạo/sửa KPI (PATCH .../strategy/kpis/:kpiId)",
            "",
            "Khi nào dùng:",
            "  • Kickoff client · review quý",
            "  • Link initiative vào goal cụ thể",
            "  • Báo cáo lãnh đạo từ OKR tree",
        ],
    )

    slide = prs.slides.add_slide(_blank(prs))
    _header(slide, "S-06 Research Console", "/seo/research · Flow F1")
    _bullets(
        slide,
        [
            "7 tabs: Keywords · Questions · Entities · Clusters · SERP · Pages · Opportunities",
            "",
            "Quy trình điển hình:",
            "  1. Import keyword CSV",
            "  2. Gom cluster + assign keyword",
            "  3. Capture SERP (PTT_SERP_PROVIDER stub/live)",
            "  4. Sync GSC pages",
            "  5. Auto-link entities ↔ clusters",
            "  6. Chọn keyword → Tạo brief → Preview → Tạo content",
            "",
            "Output: card Brief Ready trên Content Pipeline",
        ],
    )

    slide = prs.slides.add_slide(_blank(prs))
    _header(slide, "S-07 Content Pipeline", "/seo/content · 13 giai đoạn")
    _bullets(
        slide,
        [
            "Idea → Researching → Brief Ready → In Writing",
            "→ SEO Review → AEO Review → Technical Review",
            "→ Client Review → Approved → Published",
            "→ Monitoring → Refresh Required → Archived",
            "",
            "P2 filters: Full · Review only · Cần refresh",
            "Review kanban: SEO/AEO review tách biệt",
            "",
            "Ai làm gì:",
            "  Writer: In Writing → submit review",
            "  SEO: SEO Review approve/reject",
            "  AEO: AEO Review + checklist",
            "  AM: Client Review (hoặc portal client)",
        ],
        size=12,
    )

    slide = prs.slides.add_slide(_blank(prs))
    _header(slide, "S-08 Publish & Governance", "Content detail /seo/content/:id")
    _bullets(
        slide,
        [
            "Governance (PTT_SEO_GOVERNANCE_ENABLED=1):",
            "  Chặn publish nếu thiếu meta title, schema lỗi, chưa QA",
            "",
            "Publish CMS:",
            "  • Thủ công: Approved → Publish → CMS",
            "  • Tự động: PTT_SEO_CMS_AUTO_PUBLISH=1 → webhook khi Published",
            "",
            "Audit trail: lịch sử approve/reject trên content detail",
            "Runbook: runbooks/seo-cms-webhook-pilot.md",
        ],
    )

    slide = prs.slides.add_slide(_blank(prs))
    _header(slide, "S-09 Technical Console", "/seo/technical")
    _bullets(
        slide,
        [
            "Issue backlog: import crawl CSV → triage → gán fix → task CRM",
            "GSC OAuth: Kết nối Google → sync clicks/impressions hàng ngày",
            "GA4 OAuth: sessions · conversions · revenue (attribution S-12)",
            "CWV panel: PageSpeed ingest (PTT_CWV_STUB staging)",
            "Crawl connector: POST webhook ingest lịch (Gate E2)",
            "",
            "Onboard bắt buộc mỗi client pilot:",
            "  Technical → OAuth GSC + GA4 → Sync → verify clicks > 0",
        ],
    )

    slide = prs.slides.add_slide(_blank(prs))
    _header(slide, "S-10 AEO Console", "/seo/aeo")
    _bullets(
        slide,
        [
            "Question bank + coverage map theo client",
            "Batch scan: enqueue job quét AI visibility (Anthropic)",
            "Readiness score · mention trends · gap notes",
            "",
            "Luồng:",
            "  Research questions → AEO bank → Scan → Gap → FAQ content",
            "Gate D: PTT_AEO_SCHEDULE_ENABLED weekly auto-draft",
            "",
            "P2 a11y: aria-live polite cho kết quả scan",
        ],
    )

    slide = prs.slides.add_slide(_blank(prs))
    _header(slide, "S-12 Reporting & Attribution", "/seo/reports")
    _bullets(
        slide,
        [
            "Dashboard types: Executive · GSC · Content · Technical · Ops · BI",
            "Charts: GSC sparkline · severity bar (crm_seo_charts pattern)",
            "",
            "P1 Organic attribution panel:",
            "  Sessions · conversions · revenue · top landing pages",
            "  Điều kiện: GA4 sync + revenue events",
            "",
            "Export PDF (cap reports) · ClickHouse BI",
            "Lịch email: weekly/monthly cadence · Gửi ngay test",
        ],
    )

    _table_slide(
        prs,
        "Module bổ sung",
        ["Module", "Route", "Ghi chú"],
        [
            ["Automations", "/seo/automations", "Alerts · Run checks · Slack/Teams"],
            ["Governance", "/seo/governance", "Policy engine · SOP link"],
            ["Rank / SOV", "/seo/ranks", "Capture SERP · PTT_RANK_LIVE_ENABLED"],
            ["Freshness", "/seo/freshness", "Decay queue · weekly scan"],
            ["Experiments", "/seo/experiments", "Flag PTT_SEO_EXPERIMENTS_ENABLED"],
            ["BI / CMS", "/seo/bi · /seo/cms", "Grafana · webhook pilot"],
        ],
        col_widths=[1.2, 1.5, 3.0],
    )

    slide = prs.slides.add_slide(_blank(prs))
    _header(slide, "Luồng vận hành hàng ngày", "SEO Strategist / Writer / Tech SEO")
    _bullets(
        slide,
        [
            "Sáng — sau cron GSC/GA4 (06:xx):",
            "  Hub → sync OK? → client health → critical issues",
            "",
            "Strategist:",
            "  Research opportunities → brief → assign writer",
            "",
            "Writer:",
            "  Pipeline In Writing → submit SEO/AEO review",
            "",
            "Tech SEO:",
            "  Import crawl · resolve issues · CWV watch",
            "",
            "Cuối tuần:",
            "  Freshness queue · Rank capture · báo cáo client",
        ],
        size=12,
    )

    _table_slide(
        prs,
        "Onboard client SEO mới — Checklist",
        ["Bước", "Màn hình", "Done?"],
        [
            ["1", "CRM: tạo/chọn customer", "☐"],
            ["2", "/seo/clients/:id Settings — domain, tier", "☐"],
            ["3", "/seo/technical — GSC + GA4 OAuth", "☐"],
            ["4", "Sync + verify GSC clicks > 0", "☐"],
            ["5", "/seo/research — import keywords", "☐"],
            ["6", "/seo/strategy — OKR (nếu Gate E)", "☐"],
            ["7", "Hub hiển thị KPI T+1", "☐"],
            ["8", "Portal map (tuỳ chọn)", "☐"],
        ],
        subtitle="In A4: docs/forms/seo-aeo-ops-checklist-a4.html",
        col_widths=[0.5, 3.5, 0.8],
    )

    _table_slide(
        prs,
        "Xử lý sự cố nhanh",
        ["Triệu chứng", "Kiểm tra", "Fix"],
        [
            ["Menu SEO ẩn", "Cap crm_seo_aeo + NEXT_PUBLIC flags", "Admin phân quyền · rebuild ops-web"],
            ["GSC 0 rows", "OAuth · site URL", "Technical reconnect"],
            ["Attribution trống", "GA4 revenue sync", "Sync + cron daily"],
            ["Publish fail", "Governance + CMS secret", "Settings test webhook"],
            ["SERP stub", "PTT_SERP_PROVIDER", "SerpAPI key live"],
        ],
        col_widths=[1.5, 2.0, 2.8],
    )

    slide = prs.slides.add_slide(_blank(prs))
    _header(slide, "Go-live Gate A", "/seo/gate-a")
    _bullets(
        slide,
        [
            "Soak ≥7 ngày: ./scripts/phase5_soak_record.sh",
            "QA automated:",
            "  ./scripts/seo_handoff_gate.sh",
            "  SEO_HANDOFF_SKIP_E2E=0 ./scripts/playwright_ops_seo_handoff_e2e.sh",
            "  ./scripts/seo_gate_a_cutover_gate.sh",
            "",
            "Staged prod flags:",
            "  1. PTT_SEO_GOVERNANCE_ENABLED=1",
            "  2. PTT_PORTAL_SEO_ENABLED=1 (sau portal map)",
            "  3. PTT_SEO_EXPERIMENTS_ENABLED=1 (sau UAT)",
            "",
            "Sign-off: runbooks/phase5-prod-signoff-checklist.md",
        ],
        size=12,
    )

    slide = prs.slides.add_slide(_blank(prs))
    _header(slide, "Tài liệu & tài nguyên đào tạo")
    _bullets(
        slide,
        [
            "Hướng dẫn đầy đủ: docs/huong-dan-seo-aeo-ops.md (v2.0)",
            "Master spec: docs/SPEC_SEO_AEO_OPERATING_SYSTEM.md",
            "UI/UX: docs/SPEC_UI_UX_SEO_AEO.md",
            "Roadmap: docs/SEO_AEO_COMPLETION_ROADMAP.md",
            "",
            "Checklist in A4: docs/forms/seo-aeo-ops-checklist-a4.html",
            "Tạo lại slide: python3 scripts/generate_seo_aeo_training_pptx.py",
            "",
            "Env pilot: deploy/env.seo-aeo-pilot.example",
            "Gate A prod: deploy/env.seo-gate-a-prod.example",
        ],
    )

    _closing_slide(prs)

    OUTPUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUTPUT))
    return OUTPUT


def main() -> int:
    path = build()
    print(f"Wrote {path} ({path.stat().st_size // 1024} KB)")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
